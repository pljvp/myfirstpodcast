"""
AI Podcast Pipeline v3.0 - Enhanced Debug Mode
- Better error logging
- Shows WHY retry is happening
- Verbose mode for troubleshooting
"""
import os
import json
import subprocess
import re
import time
from pathlib import Path
from datetime import datetime
from dotenv import load_dotenv
from anthropic import Anthropic
import requests

# TTS Provider modules
from providers import ElevenLabsProvider, CartesiaProvider, substitute_template_placeholders

# Document reading libraries (optional - graceful fallback if not installed)
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Load environment variables
config_path = Path(__file__).parent / 'config' / '.env'
load_dotenv(config_path)

# Debug mode - set to True for verbose logging
DEBUG_VERBOSE = True




def get_text_editor():
    """Get appropriate text editor for current OS"""
    import platform
    
    # Check environment variable first
    if 'EDITOR' in os.environ:
        return os.environ['EDITOR']
    
    # Detect OS
    system = platform.system()
    
    if system == 'Windows':
        return 'notepad'
    else:  # Linux, Unix
        return 'nano'

def log_debug(message):
    """Print debug message if verbose mode enabled"""
    if DEBUG_VERBOSE:
        print(f"[VERBOSE] {message}")


def load_config():
    """Load podcast_config.json from config folder"""
    config_file = Path(__file__).parent / 'config' / 'podcast_config.json'
    with open(config_file, 'r', encoding='utf-8') as f:
        return json.load(f)


def create_project_structure(project_name):
    """Create project folder with subdirectories"""
    base_path = Path(f"./projects/{project_name}")
    (base_path / "prompts").mkdir(parents=True, exist_ok=True)
    (base_path / "sources").mkdir(parents=True, exist_ok=True)
    (base_path / "scripts").mkdir(parents=True, exist_ok=True)
    (base_path / "audio").mkdir(parents=True, exist_ok=True)
    (base_path / "debug").mkdir(parents=True, exist_ok=True)
    
    sources_file = base_path / "sources" / "sources_list.txt"
    if not sources_file.exists():
        with open(sources_file, 'w', encoding='utf-8') as f:
            f.write(f"Research Sources for {project_name}\n\n")
            f.write("Primary Sources:\n- \n\n")
            f.write("Background Reading:\n- \n\n")
            f.write("Key Points to Cover:\n- \n")
    
    context_file = base_path / "sources" / "research_context.txt"
    if not context_file.exists():
        # Check if there's a default template to use
        default_template = Path("templates/research_contexts/default.txt")
        
        if default_template.exists():
            log_debug(f"Using default research context template: {default_template}")
            with open(default_template, 'r', encoding='utf-8') as f:
                template_content = f.read()
            with open(context_file, 'w', encoding='utf-8') as f:
                f.write(template_content.replace("{project_name}", project_name))
            print(f"  ✓ Using default research context template")
        else:
            # Create minimal default
            log_debug("No template found, creating minimal default")
            with open(context_file, 'w', encoding='utf-8') as f:
                f.write(f"Research Context for {project_name}\n\n")
                f.write("=== RESEARCH INSTRUCTIONS ===\n")
                f.write("Number of sources to find: 5-10\n")
                f.write("Focus on recent (2024-2025) developments\n\n")
                f.write("=== CONTEXT AND FOCUS AREAS ===\n")
                f.write("(Describe what Claude should focus on during research)\n\n")
                f.write("=== SPECIFIC QUESTIONS TO ANSWER ===\n")
                f.write("1. What are the latest developments?\n")
                f.write("2. What are the practical applications?\n")
                f.write("3. What are experts saying?\n\n")
                f.write("=== AUDIENCE CONSIDERATIONS ===\n")
                f.write("Intelligent general audience - explain jargon, use analogies\n")
    else:
        print(f"  ✓ Using existing research context (project-specific)")
    
    return base_path


def load_template(template_path, variables):
    """Load template and substitute variables"""
    with open(template_path, 'r', encoding='utf-8') as f:
        template = f.read()
    for key, value in variables.items():
        template = template.replace(f"{{{key}}}", str(value))
    return template



def build_scenario_context(scenario_type, topic_description, language='german'):
    """Build scenario-specific context instructions for test templates"""
    
    # German contexts
    contexts_de = {
        'road': f"""SZENARIO-KONTEXT:

SITUATION: Paar im Auto, verloren und streitet über Navigation

ORT: {topic_description}

EMOTIONALER BOGEN:
- Anfang: Selbstsicher, aber beginnende Zweifel
- Mitte: Zunehmende Frustration, spielerisches Gezänk
- Höhepunkt: Erkenntnis, dass sie völlig verloren sind
- Auflösung: Humor und Akzeptanz

SCHLÜSSEL-ELEMENTE:
- GPS gibt falsche Anweisungen
- Uneinigkeit über "links" vs "rechts"
- Verweis auf vergangene Navigations-Desaster
- Physische Komik (Beinahe-Unfälle, falsche Abbiegungen)
- Eine Person verteidigt ihre Navigations-Fähigkeiten
- Andere Person weist skeptisch auf Fehler hin

TON: Frustriert aber liebevoll, komische Spannung""",
        
        'cook': f"""SZENARIO-KONTEXT:

SITUATION: Zwei Personen kochen zusammen, Desaster entfaltet sich

REZEPT: {topic_description}

EMOTIONALER BOGEN:
- Anfang: Optimistisch, lockerer Ansatz
- Mitte: Dinge gehen schief, Panik steigt
- Höhepunkt: Rauchmelder / komplettes Versagen
- Auflösung: Essen bestellen, darüber lachen

SCHLÜSSEL-ELEMENTE:
- Eine Person improvisiert, andere folgt Rezept
- Mess-Desaster (zu viel/wenig Zutaten)
- Temperatur-Probleme (zu heiß/kalt)
- Rauch/Brandgeruch
- Uneinigkeit über "Kochen ist Kunst vs Wissenschaft"
- Physische Reaktionen (gasps, Panik)

TON: Eskalierende Chaos, erhalten Zuneigung""",
        
        'mvie': f"""SZENARIO-KONTEXT:

SITUATION: Film-Enthusiasten debattieren über berühmte Szene

FILM/SZENE: {topic_description}

EMOTIONALER BOGEN:
- Anfang: Eine Person leidenschaftlich, andere neugierig
- Mitte: Analytische Diskussion, Meinungsverschiedenheiten
- Höhepunkt: Offenbarung oder lustige Beobachtung
- Auflösung: Einigung oder agree-to-disagree

SCHLÜSSEL-ELEMENTE:
- Zitat oder Verweis auf tatsächliche Szene
- Debatte über Interpretation/Bedeutung
- Eine Person analytisch, andere anfangs abweisend
- Entdeckung neuer Perspektive
- Pop-Kultur-Referenzen
- Wechsel von Skepsis zu Wertschätzung (oder umgekehrt)

TON: Intellektuell aber spielerisch, leidenschaftliche Diskussion"""
    }
    
    # English contexts
    contexts_en = {
        'road': f"""SCENARIO CONTEXT:

SITUATION: Couple in car, lost and arguing about navigation

LOCATION: {topic_description}

EMOTIONAL ARC:
- Start: Confident but starting to doubt
- Middle: Escalating frustration, playful bickering
- Climax: Realization they're completely lost
- Resolution: Humor and acceptance

KEY ELEMENTS:
- GPS giving wrong directions
- Disagreement about "left" vs "right"
- Reference to past navigation disasters
- Physical comedy (near-misses, wrong turns)
- One person defending their navigation skills
- Other person skeptically pointing out mistakes

TONE: Frustrated but affectionate, comedic tension""",
        
        'cook': f"""SCENARIO CONTEXT:

SITUATION: Two people cooking together, disaster unfolds

RECIPE: {topic_description}

EMOTIONAL ARC:
- Start: Optimistic, casual approach
- Middle: Things going wrong, panic rising
- Climax: Smoke alarm / complete failure
- Resolution: Ordering takeout, laughing about it

KEY ELEMENTS:
- One person improvising, other following recipe
- Measurement disasters
- Temperature problems
- Smoke/burning smell
- Disagreement about "cooking is art vs science"
- Physical reactions (gasps, panic)

TONE: Escalating chaos, maintained affection""",
        
        'mvie': f"""SCENARIO CONTEXT:

SITUATION: Film enthusiasts debating famous scene

MOVIE/SCENE: {topic_description}

EMOTIONAL ARC:
- Start: One passionate, other curious
- Middle: Analytical discussion, disagreements
- Climax: Revelation or funny observation
- Resolution: Agreement or agree-to-disagree

KEY ELEMENTS:
- Quotation or reference to actual scene
- Debate about interpretation/meaning
- One analytical, other dismissive initially
- Discovery of new perspective
- Pop culture references
- Shift from skepticism to appreciation (or vice versa)

TONE: Intellectual but playful, passionate discussion"""
    }
    
    # Dutch contexts
    contexts_nl = {
        'road': f"""SCENARIO CONTEXT:

SITUATIE: Stel in auto, verdwaald en ruzie over navigatie

LOCATIE: {topic_description}

EMOTIONELE BOOG:
- Begin: Zelfverzekerd maar beginnende twijfel
- Midden: Toenemende frustratie, speelse ruzie
- Hoogtepunt: Besef dat ze helemaal verdwaald zijn
- Oplossing: Humor en acceptatie

SLEUTELELEMENTEN:
- GPS geeft verkeerde aanwijzingen
- Onenigheid over "links" vs "rechts"
- Verwijzing naar eerdere navigatie-rampen
- Fysieke komedie (bijna-ongelukken, verkeerde afslag)
- Eén persoon verdedigt navigatie-vaardigheden
- Ander persoon wijst skeptisch op fouten

TOON: Gefrustreerd maar liefdevol, komische spanning""",
        
        'cook': f"""SCENARIO CONTEXT:

SITUATIE: Twee mensen samen koken, ramp ontvouwt zich

RECEPT: {topic_description}

EMOTIONELE BOOG:
- Begin: Optimistisch, relaxte aanpak
- Midden: Dingen gaan mis, paniek stijgt
- Hoogtepunt: Rookalarm / compleet falen
- Oplossing: Eten bestellen, erom lachen

SLEUTELELEMENTEN:
- Eén persoon improviseert, ander volgt recept
- Meet-rampen (te veel/weinig ingrediënten)
- Temperatuur problemen (te heet/koud)
- Rook/brandlucht
- Onenigheid over "koken is kunst vs wetenschap"
- Fysieke reacties (gasps, paniek)

TOON: Escalerende chaos, behouden genegenheid""",
        
        'mvie': f"""SCENARIO CONTEXT:

SITUATIE: Film-enthousiastelingen debatteren over beroemde scène

FILM/SCÈNE: {topic_description}

EMOTIONELE BOOG:
- Begin: Eén gepassioneerd, ander nieuwsgierig
- Midden: Analytische discussie, meningsverschillen
- Hoogtepunt: Onthulling of grappige observatie
- Oplossing: Akkoord of agree-to-disagree

SLEUTELELEMENTEN:
- Citaat of verwijzing naar werkelijke scène
- Debat over interpretatie/betekenis
- Eén analytisch, ander aanvankelijk afwijzend
- Ontdekking van nieuw perspectief
- Popcultuur referenties
- Verschuiving van scepsis naar waardering (of omgekeerd)

TOON: Intellectueel maar speels, gepassioneerde discussie"""
    }
    
    # Select context based on language
    if language == 'german':
        contexts = contexts_de
    elif language == 'english':
        contexts = contexts_en
    elif language == 'dutch':
        contexts = contexts_nl
    else:
        contexts = contexts_en
    
    return contexts.get(scenario_type, '')


def save_script_test(script, project_name, language_code, topic_tag, provider_tag, draft_number):
    """Save test script with scenario tag: test_LANG_DATE_scenario-topic_PROV_draftN.txt"""
    timestamp = datetime.now().strftime('%Y-%m-%d_%H-%M')
    lang_upper = language_code.upper()
    
    # test_DE_2025-11-29_19-30_road-prsd_CRTS_draft1.txt
    filename = f"{project_name.lower()}_{lang_upper}_{timestamp}_{topic_tag}_{provider_tag}_draft{draft_number}.txt"
    
    path = Path(f"./projects/{project_name}/scripts/{filename}")
    with open(path, 'w', encoding='utf-8') as f:
        f.write(script)
    
    return path


def get_user_input(prompt, options=None):
    """Get user input with optional menu"""
    if options:
        print(f"\n{prompt}")
        for i, option in enumerate(options, 1):
            print(f"    {i}. {option}")
        while True:
            try:
                choice = int(input("Choice: "))
                if 1 <= choice <= len(options):
                    return choice - 1
                print(f"Please enter a number between 1 and {len(options)}")
            except ValueError:
                print("Please enter a valid number")
    else:
        return input(f"{prompt}: ")


def generate_script(prompt, api_key):
    """Call Claude API with prompt"""
    print("\n" + "="*60)
    print("CLAUDE IS WORKING...")
    print("="*60)
    print("- Conducting online research")
    print("- Analyzing sources")
    print("- Generating podcast script")
    print("- Formatting dialogue")
    print("")
    print("This may take 30-60 seconds...")
    print("="*60 + "\n")
    
    client = Anthropic(api_key=api_key)
    
    try:
        response = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=8000,
            messages=[{"role": "user", "content": prompt}]
        )
        print("✓ Script generated successfully!\n")
        
        # Track usage
        usage = response.usage
        print(f"[USAGE] Claude - Input: {usage.input_tokens} tokens, Output: {usage.output_tokens} tokens")
        
        return response.content[0].text, usage
    except Exception as e:
        print(f"\n✗ Error calling Claude API: {e}\n")
        return None, None


def revise_script(original_script, revision_guidance, api_key):
    """Request Claude to revise script"""
    prompt = f"""Here is a podcast script:

{original_script}

Please revise this script according to the following guidance:
{revision_guidance}

Provide the complete revised script maintaining the same format with Speaker A and Speaker B labels."""

    print("\n" + "="*60)
    print("CLAUDE IS REVISING SCRIPT...")
    print("="*60)
    print("- Analyzing your feedback")
    print("- Updating script content")
    print("- Maintaining dialogue format")
    print("")
    print("This may take 30-60 seconds...")
    print("="*60 + "\n")

    return generate_script(prompt, api_key)


# =============================================================================
# MULTI-CALL ARCHITECTURE FUNCTIONS
# =============================================================================

import math

def estimate_api_calls(duration_minutes, doc_count, web_source_count, config):
    """
    Calculate API calls needed and estimated cost.
    Returns dict with breakdown and totals.
    """
    gen_config = config.get('script_generation', {})
    words_per_call = gen_config.get('words_per_call', 2000)
    docs_per_batch = gen_config.get('docs_per_batch', 3)
    sources_per_call = gen_config.get('sources_per_research_call', 10)

    word_count = duration_minutes * 222

    # Research calls (sources_per_call sources per call)
    research_calls = math.ceil(web_source_count / sources_per_call) if web_source_count > 0 else 0

    # Document processing calls (docs_per_batch docs per batch)
    doc_calls = math.ceil(doc_count / docs_per_batch) if doc_count > 0 else 0

    # Outline call (always 1 for multi-call mode)
    outline_calls = 1

    # Script generation calls (~words_per_call words per call)
    script_calls = math.ceil(word_count / words_per_call)

    # Lightweight synthesis: 1 small call per transition (N-1 for N sections)
    synthesis_calls = script_calls - 1 if script_calls > 1 else 0

    total_calls = research_calls + doc_calls + outline_calls + script_calls + synthesis_calls

    # Cost estimation (Sonnet: ~$0.003/1K input + $0.015/1K output)
    # Script calls: ~3K input, ~2K output = ~$0.04 each
    # Synthesis calls: ~500 input, ~500 output = ~$0.01 each (lightweight)
    script_cost = script_calls * 0.04
    synthesis_cost = synthesis_calls * 0.01
    other_cost = (research_calls + doc_calls + outline_calls) * 0.04
    estimated_cost = script_cost + synthesis_cost + other_cost

    return {
        'word_count': word_count,
        'research_calls': research_calls,
        'doc_calls': doc_calls,
        'outline_calls': outline_calls,
        'script_calls': script_calls,
        'synthesis_calls': synthesis_calls,
        'total_calls': total_calls,
        'estimated_cost': estimated_cost
    }


def display_generation_plan(duration, doc_count, web_source_count, config):
    """
    Display the generation plan and get user confirmation.
    Returns True if user confirms, False otherwise.
    """
    estimate = estimate_api_calls(duration, doc_count, web_source_count, config)

    print("\n" + "="*60)
    print("GENERATION PLAN")
    print("="*60)
    print(f"  Duration: {duration} minutes (~{estimate['word_count']} words)")
    print(f"  Web sources: {web_source_count} requested")
    print(f"  Documents: {doc_count} file(s)")
    print("-"*60)
    print("  ESTIMATED CLAUDE API CALLS:")
    if estimate['research_calls'] > 0:
        print(f"  ├─ Research phase:     {estimate['research_calls']} call(s)")
    if estimate['doc_calls'] > 0:
        print(f"  ├─ Document summaries: {estimate['doc_calls']} call(s)")
    print(f"  ├─ Outline generation: {estimate['outline_calls']} call")
    print(f"  ├─ Script generation:  {estimate['script_calls']} call(s)")
    if estimate['synthesis_calls'] > 0:
        print(f"  └─ Transition smoothing: {estimate['synthesis_calls']} call(s) (lightweight)")
    else:
        print(f"  └─ Transition smoothing: 0 (single section)")
    print("-"*60)
    print(f"  TOTAL: {estimate['total_calls']} Claude API calls")
    print(f"  Est. cost: ~${estimate['estimated_cost']:.2f}")
    print("="*60)

    confirm = input("\nProceed with generation? (Y/n): ").strip().lower()
    return confirm != 'n'


def research_web_sources(topic, research_context, source_count, api_key, config):
    """
    Conduct web research in multiple calls if needed.
    Returns combined research findings.
    """
    gen_config = config.get('script_generation', {})
    sources_per_call = gen_config.get('sources_per_research_call', 10)
    show_progress = gen_config.get('show_progress', True)

    client = Anthropic(api_key=api_key)

    total_calls = math.ceil(source_count / sources_per_call)
    all_findings = []
    all_sources = []

    for call_num in range(1, total_calls + 1):
        sources_this_call = min(sources_per_call, source_count - (call_num - 1) * sources_per_call)

        if show_progress:
            print(f"\n[RESEARCH] Call {call_num}/{total_calls}: Finding {sources_this_call} sources...")

        prompt = f"""You are a research assistant. Search the web and find {sources_this_call} high-quality, recent sources about:

TOPIC: {topic}

RESEARCH FOCUS:
{research_context}

For each source:
1. Search for recent (2024-2025 preferred) authoritative content
2. Extract key insights, facts, statistics, and expert opinions
3. Note any controversies or different perspectives

OUTPUT FORMAT:
### SOURCE 1: [Title]
URL: [url]
KEY INSIGHTS:
- [insight 1]
- [insight 2]
- [insight 3]

### SOURCE 2: [Title]
...

After all sources, provide:
### SYNTHESIS
[2-3 paragraph summary of the most important findings across all sources]
"""

        try:
            response = client.messages.create(
                model="claude-sonnet-4-20250514",
                max_tokens=4000,
                messages=[{"role": "user", "content": prompt}]
            )

            findings = response.content[0].text
            all_findings.append(findings)

            usage = response.usage
            if show_progress:
                print(f"    ✓ Found sources (Input: {usage.input_tokens}, Output: {usage.output_tokens} tokens)")

        except Exception as e:
            print(f"    ✗ Research call {call_num} failed: {e}")
            continue

    # Combine all findings
    if show_progress and all_findings:
        print(f"\n[Finalizing] Compiling research summary...")
    combined = "\n\n---\n\n".join(all_findings)
    return combined


def process_documents_batched(documents_text, project_name, api_key, config):
    """
    Process source documents in batches, creating summaries.
    Returns combined document summaries.
    """
    gen_config = config.get('script_generation', {})
    docs_per_batch = gen_config.get('docs_per_batch', 3)
    show_progress = gen_config.get('show_progress', True)

    # Split documents by the ### SOURCE: marker
    doc_sections = documents_text.split('### SOURCE:')
    doc_sections = [d.strip() for d in doc_sections if d.strip()]

    if not doc_sections:
        return ""

    client = Anthropic(api_key=api_key)

    # Create batches
    batches = []
    for i in range(0, len(doc_sections), docs_per_batch):
        batch = doc_sections[i:i + docs_per_batch]
        batches.append(batch)

    all_summaries = []

    for batch_num, batch in enumerate(batches, 1):
        if show_progress:
            print(f"\n[DOCUMENTS] Processing batch {batch_num}/{len(batches)} ({len(batch)} docs)...")

        batch_text = "\n\n### SOURCE:".join(batch)

        prompt = f"""Summarize the following source documents for use in a podcast script.

For each document, extract:
1. Main thesis/argument
2. Key facts, statistics, and data points
3. Notable quotes or expert opinions
4. Practical examples or case studies

DOCUMENTS:
### SOURCE:{batch_text}

OUTPUT FORMAT:
### DOCUMENT SUMMARY 1: [filename]
MAIN POINTS:
- [point 1]
- [point 2]
KEY DATA:
- [statistic or fact]
USABLE QUOTES:
- "[quote]"

### DOCUMENT SUMMARY 2: [filename]
...

Keep summaries concise but preserve specific details that would be valuable for podcast discussion.
"""

        try:
            response = client.messages.create(
                model="claude-sonnet-4-20250514",
                max_tokens=3000,
                messages=[{"role": "user", "content": prompt}]
            )

            summary = response.content[0].text
            all_summaries.append(summary)

            usage = response.usage
            if show_progress:
                print(f"    ✓ Batch summarized (Input: {usage.input_tokens}, Output: {usage.output_tokens} tokens)")

        except Exception as e:
            print(f"    ✗ Document batch {batch_num} failed: {e}")
            continue

    if show_progress and all_summaries:
        print(f"\n[Finalizing] Combining document summaries...")
    return "\n\n".join(all_summaries)


def generate_outline(topic, duration, word_count, research_summary, doc_summary, style_description, language, api_key, config):
    """
    Generate a structured outline for the podcast.
    Returns outline text that guides script generation.
    """
    show_progress = config.get('script_generation', {}).get('show_progress', True)

    if show_progress:
        print(f"\n[OUTLINE] Generating story arc and section breakdown...")

    client = Anthropic(api_key=api_key)

    # Calculate sections based on word count
    words_per_call = config.get('script_generation', {}).get('words_per_call', 2000)
    overshoot_factor = config.get('script_generation', {}).get('overshoot_factor', 1.5)
    num_sections = math.ceil(word_count / words_per_call)
    words_per_section = int((word_count / num_sections) * overshoot_factor)
    total_target_words = words_per_section * num_sections

    prompt = f"""Create a detailed podcast outline for a {duration}-minute episode (~{total_target_words} words total).

TOPIC: {topic}
STYLE: {style_description}
LANGUAGE: {language}
NUMBER OF SECTIONS: {num_sections} (each ~{words_per_section} words)

RESEARCH FINDINGS:
{research_summary[:8000] if research_summary else "No web research provided."}

DOCUMENT INSIGHTS:
{doc_summary[:4000] if doc_summary else "No source documents provided."}

CREATE AN OUTLINE WITH:

1. **OVERALL ARC**: Describe the narrative journey (hook → exploration → insight → conclusion)

2. **SPEAKER DYNAMICS**:
   - Speaker A (LEAD VOICE - female): knowledgeable expert, enthusiastic explainer, drives the conversation
   - Speaker B (male): curious questioner, friendly skeptic, asks follow-up questions

3. **SECTION BREAKDOWN** (one for each of the {num_sections} sections):

### SECTION 1: [Title] (~{words_per_section} words)
SPEAKER LEAD: A (Speaker A should lead most sections)
CONTENT:
- Opening hook: [specific hook idea]
- Key points to cover: [bullet list]
- Facts/stats to include: [from research]
- Emotional beats: [curiosity, surprise, humor, etc.]
TRANSITION TO NEXT: "[Exact transition phrase]"

### SECTION 2: [Title] (~{words_per_section} words)
...

4. **KEY MOMENTS**: List 3-5 specific moments that should feel memorable (a surprising fact, a funny exchange, an "aha" moment)

5. **CLOSING**: How the episode should end (call to action, reflection, teaser)

Be specific. Include actual facts from the research. This outline will guide multiple script-generation calls.
"""

    try:
        response = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=4000,
            messages=[{"role": "user", "content": prompt}]
        )

        outline = response.content[0].text
        usage = response.usage

        if show_progress:
            print(f"    ✓ Outline created (Input: {usage.input_tokens}, Output: {usage.output_tokens} tokens)")

        return outline, usage

    except Exception as e:
        print(f"    ✗ Outline generation failed: {e}")
        return None, None


def generate_script_section(section_num, total_sections, outline, previous_section_end,
                           target_words, style_template, language, api_key, config, provider='elevenlabs'):
    """
    Generate a single section of the podcast script.
    Returns the section text.

    Args:
        provider: 'elevenlabs' or 'cartesia' - determines emotion tag instructions
    """
    show_progress = config.get('script_generation', {}).get('show_progress', True)

    if show_progress:
        print(f"\n[SCRIPT] Generating section {section_num}/{total_sections} (~{target_words} words)...")

    client = Anthropic(api_key=api_key)

    # Context from previous section for continuity
    continuity_context = ""
    if previous_section_end:
        continuity_context = f"""
PREVIOUS SECTION ENDING (maintain continuity):
{previous_section_end}

Continue naturally from this point. Do NOT repeat content.
"""

    # Section-specific instructions
    if section_num == 1:
        section_instruction = "This is the OPENING section. Start with an engaging hook. Introduce the topic and speakers' dynamic."
    elif section_num == total_sections:
        section_instruction = "This is the CLOSING section. Build toward a satisfying conclusion. Include summary and call-to-action."
    else:
        section_instruction = f"This is the MIDDLE section {section_num}. Continue building on previous content. Maintain energy and introduce new angles."

    prompt = f"""Generate section {section_num} of {total_sections} for a podcast script.

MINIMUM: {target_words} words for this section (do not write less)
LANGUAGE: {language}

{section_instruction}

OUTLINE (follow this structure):
{outline}

{continuity_context}

STYLE REQUIREMENTS:
{style_template[:2000] if style_template else "Natural, conversational dialogue between Speaker A and Speaker B."}

CRITICAL FORMAT REQUIREMENTS:
1. Start IMMEDIATELY with dialogue - NO title, NO header, NO introduction text
2. Use EXACTLY this format: "Speaker A:" or "Speaker B:" (NO asterisks, NO bold, NO markdown)
3. NO blank lines between dialogue segments - each line follows immediately after the previous
4. Place emotion tags at the START of each line, AFTER "Speaker A:" or "Speaker B:"
5. If emotion changes mid-thought, START A NEW LINE with the new speaker label and emotion
6. End at a natural transition point (NOT mid-sentence)
7. Do NOT include sources, titles, dividers (---), or any non-dialogue content

CORRECT FORMAT EXAMPLE:
Speaker A: [excited] This is amazing news!
Speaker B: [curious] Tell me more about it.
Speaker A: [thoughtful] Well, the research shows...
Speaker A: [surprised] And this is the fascinating part!

WRONG FORMAT (DO NOT USE):
Speaker A: [excited] This is amazing! [thoughtful] But we need to consider...
(Mid-line emotion changes are NOT allowed - start a new line instead)

**Speaker A:** [excited] This is amazing news!
(Asterisks/bold are NOT allowed)

Generate the script section now (start directly with "Speaker A:" or "Speaker B:"):
"""

    # Calculate max_tokens based on target words (1.5 tokens/word + buffer)
    section_max_tokens = int(target_words * 1.5) + 500

    try:
        response = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=section_max_tokens,
            messages=[{"role": "user", "content": prompt}]
        )

        section_text = response.content[0].text
        usage = response.usage

        if show_progress:
            word_count = len(section_text.split())
            print(f"    ✓ Section {section_num} generated ({word_count} words, Output: {usage.output_tokens} tokens)")

        return section_text, usage

    except Exception as e:
        print(f"    ✗ Section {section_num} generation failed: {e}")
        return None, None


def generate_script_multi_call(topic, duration, word_count, outline, style_template, language, api_key, config, provider='elevenlabs'):
    """
    Orchestrate multi-call script generation using the outline.
    Returns complete script text.

    Args:
        provider: 'elevenlabs' or 'cartesia' - determines emotion tag instructions
    """
    gen_config = config.get('script_generation', {})
    words_per_call = gen_config.get('words_per_call', 2000)
    # LLMs typically generate 65-70% of requested words - compensate with overshoot
    overshoot_factor = gen_config.get('overshoot_factor', 1.4)

    num_sections = math.ceil(word_count / words_per_call)
    words_per_section = int((word_count / num_sections) * overshoot_factor)

    print("\n" + "="*60)
    print(f"GENERATING SCRIPT ({num_sections} sections)")
    print("="*60)

    sections = []
    previous_ending = None
    total_usage = {'input': 0, 'output': 0}

    for section_num in range(1, num_sections + 1):
        section_text, usage = generate_script_section(
            section_num=section_num,
            total_sections=num_sections,
            outline=outline,
            previous_section_end=previous_ending,
            target_words=words_per_section,
            style_template=style_template,
            language=language,
            api_key=api_key,
            config=config,
            provider=provider
        )

        if not section_text:
            print(f"    ✗ Failed to generate section {section_num}, aborting")
            return None, None

        sections.append(section_text)

        # Keep last ~500 words for continuity
        words = section_text.split()
        if len(words) > 100:
            previous_ending = ' '.join(words[-100:])
        else:
            previous_ending = section_text

        if usage:
            total_usage['input'] += usage.input_tokens
            total_usage['output'] += usage.output_tokens

    print(f"\n[SCRIPT] All {num_sections} sections generated")
    print(f"[USAGE] Total - Input: {total_usage['input']}, Output: {total_usage['output']} tokens")

    # Return sections list (not combined) for lightweight synthesis
    return sections, total_usage


def synthesize_script(raw_script, outline, language, api_key, config):
    """
    Final pass to smooth transitions and ensure consistency.
    Returns polished script.
    """
    show_progress = config.get('script_generation', {}).get('show_progress', True)

    if show_progress:
        print(f"\n[SYNTHESIS] Polishing transitions and consistency...")

    client = Anthropic(api_key=api_key)

    prompt = f"""Review and polish this podcast script. The script was generated in sections and may need smoothing.

TASKS:
1. Smooth any awkward transitions between sections
2. Ensure consistent speaker personalities throughout
3. DO NOT remove any dialogue content - preserve all material, only smooth transitions
4. Ensure emotional tags are balanced (not too many, not too few)
5. Fix any formatting inconsistencies
6. Maintain the target language: {language}

CRITICAL FORMAT REQUIREMENTS:
- Use EXACTLY "Speaker A:" and "Speaker B:" format (NO asterisks, NO bold, NO markdown like **)
- NO blank lines between dialogue segments - each line immediately follows the previous
- NO title or header at the start - begin directly with dialogue
- NO sources section, NO dividers (---), NO non-dialogue content
- Preserve all [emotion tags] in square brackets

CORRECT OUTPUT FORMAT:
Speaker A: [excited] First line of dialogue here.
Speaker B: [curious] Response follows immediately.
Speaker A: [thoughtful] And so on...

WRONG FORMAT (DO NOT USE):
**Speaker A:** [excited] Wrong format with asterisks.

**Speaker B:** [curious] Wrong with blank lines between.

OUTLINE (for reference):
{outline[:2000]}

SCRIPT TO POLISH:
{raw_script}

OUTPUT the polished script (start directly with "Speaker A:" or "Speaker B:"):
"""

    try:
        response = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=16000,
            messages=[{"role": "user", "content": prompt}]
        )

        polished = response.content[0].text
        usage = response.usage

        if show_progress:
            print(f"    ✓ Script polished (Input: {usage.input_tokens}, Output: {usage.output_tokens} tokens)")

        return polished, usage

    except Exception as e:
        print(f"    ✗ Synthesis failed: {e}")
        return raw_script, None  # Return unpolished if synthesis fails


def clean_script_format(script):
    """
    Post-process script to ensure clean format for TTS.
    Removes markdown, titles, blank lines between dialogue.
    """
    import re

    lines = script.split('\n')
    cleaned_lines = []
    in_dialogue = False

    for line in lines:
        stripped = line.strip()

        # Skip empty lines
        if not stripped:
            continue

        # Skip title lines (# headers)
        if stripped.startswith('#'):
            continue

        # Skip divider lines (---)
        if stripped.startswith('---') or stripped == '---':
            continue

        # Skip source sections
        if 'SOURCES FOUND' in stripped.upper() or 'SOURCE:' in stripped.upper():
            break  # Stop processing once we hit sources

        # Check if this is a dialogue line
        is_dialogue = any(marker in stripped.lower() for marker in
                         ['speaker a:', 'speaker b:', '**speaker a', '**speaker b'])

        if is_dialogue:
            in_dialogue = True
            # Remove markdown bold formatting
            cleaned = stripped
            cleaned = re.sub(r'\*\*Speaker ([AB]):\*\*', r'Speaker \1:', cleaned)
            cleaned = re.sub(r'\*\*Speaker ([AB])\*\*:', r'Speaker \1:', cleaned)
            cleaned = re.sub(r'\*\*(Speaker [AB]:)\*\*', r'\1', cleaned)
            # Remove any remaining ** at start
            cleaned = re.sub(r'^\*\*\s*', '', cleaned)
            cleaned_lines.append(cleaned)
        elif in_dialogue:
            # Non-dialogue line after dialogue started - might be continuation or junk
            # If it doesn't look like a header/title, include it
            if not stripped.startswith('#') and not stripped.startswith('---'):
                cleaned_lines.append(stripped)

    return '\n'.join(cleaned_lines)


def parse_script_to_segments(script):
    """
    Parse a script into dialogue segments.
    Returns list of tuples: [(speaker, full_line), ...]
    """
    segments = []
    for line in script.split('\n'):
        stripped = line.strip()
        if not stripped:
            continue

        # Check if this is a speaker line
        lower = stripped.lower()
        if lower.startswith('speaker a:') or lower.startswith('speaker b:'):
            # Determine speaker
            speaker = 'A' if 'speaker a:' in lower else 'B'
            segments.append((speaker, stripped))
        elif '**speaker a' in lower or '**speaker b' in lower:
            # Handle markdown format
            speaker = 'A' if 'speaker a' in lower else 'B'
            segments.append((speaker, stripped))

    return segments


def synthesize_single_transition(segments_before, segments_after, language, api_key):
    """
    Smooth the transition between two sections.
    Takes last N segments of section 1 and first N segments of section 2.
    Returns smoothed segments as text.
    """
    client = Anthropic(api_key=api_key)

    # Detect collision (same speaker ends section 1 and starts section 2)
    last_speaker = segments_before[-1][0] if segments_before else None
    first_speaker = segments_after[0][0] if segments_after else None
    collision = last_speaker == first_speaker

    # Build segment text
    before_text = '\n'.join([seg[1] for seg in segments_before])
    after_text = '\n'.join([seg[1] for seg in segments_after])

    collision_instruction = ""
    if collision:
        collision_instruction = f"""
COLLISION DETECTED: Section 1 ends with Speaker {last_speaker} and Section 2 starts with Speaker {last_speaker}.
You MUST fix this by either:
- Merging the two {last_speaker} segments into one natural segment
- Adding a brief bridge line from Speaker {'B' if last_speaker == 'A' else 'A'} between them
The final output MUST alternate A-B-A-B properly."""

    prompt = f"""Smooth this podcast transition between two sections.

END OF SECTION (last 4 segments):
{before_text}

START OF NEXT SECTION (first 4 segments):
{after_text}
{collision_instruction}

TASKS:
1. Create a smooth, natural transition between these segments
2. Maintain A-B-A-B alternation (fix any collision)
3. Preserve ALL factual content - do not remove information
4. Keep all [emotion tags] in square brackets
5. Maintain the language: {language}

FORMAT REQUIREMENTS:
- Use EXACTLY "Speaker A:" or "Speaker B:" format (NO asterisks, NO markdown)
- NO blank lines between segments
- Start output directly with a Speaker line

OUTPUT the smoothed transition (typically 6-10 segments):"""

    try:
        response = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=2000,
            messages=[{"role": "user", "content": prompt}]
        )
        return response.content[0].text.strip(), response.usage
    except Exception as e:
        print(f"    ✗ Transition synthesis failed: {e}")
        # Fallback: just join the segments
        return before_text + '\n' + after_text, None


def synthesize_transitions(sections, language, api_key, config):
    """
    Lightweight synthesis: only smooth the join points between sections.
    Takes list of section texts, returns combined script with smoothed transitions.
    """
    if len(sections) <= 1:
        return sections[0] if sections else "", None

    show_progress = config.get('script_generation', {}).get('show_progress', True)
    num_joins = len(sections) - 1

    if show_progress:
        print(f"\n[SYNTHESIS] Smoothing {num_joins} transition(s)...")

    total_usage = {'input': 0, 'output': 0}

    # Parse all sections into segments
    parsed_sections = [parse_script_to_segments(section) for section in sections]

    # Process each join point
    smoothed_sections = []

    for i in range(len(sections)):
        if i == 0:
            # First section: keep all but last 4 segments, they'll be in the transition
            if len(parsed_sections[i]) > 4:
                kept_segments = parsed_sections[i][:-4]
                smoothed_sections.append('\n'.join([seg[1] for seg in kept_segments]))
            # else: entire section is in transition

        if i < len(sections) - 1:
            # Get segments for this transition
            segments_before = parsed_sections[i][-4:] if len(parsed_sections[i]) >= 4 else parsed_sections[i]
            segments_after = parsed_sections[i+1][:4] if len(parsed_sections[i+1]) >= 4 else parsed_sections[i+1]

            if show_progress:
                print(f"    Smoothing transition {i+1}/{num_joins}...")

            # Synthesize this transition
            smoothed_transition, usage = synthesize_single_transition(
                segments_before=segments_before,
                segments_after=segments_after,
                language=language,
                api_key=api_key
            )

            smoothed_sections.append(smoothed_transition)

            if usage:
                total_usage['input'] += usage.input_tokens
                total_usage['output'] += usage.output_tokens
                if show_progress:
                    print(f"    ✓ Transition {i+1} smoothed ({usage.output_tokens} tokens)")

        if i == len(sections) - 1:
            # Last section: keep all but first 4 segments (they were in the transition)
            if len(parsed_sections[i]) > 4:
                kept_segments = parsed_sections[i][4:]
                smoothed_sections.append('\n'.join([seg[1] for seg in kept_segments]))
            # else: entire section was in transition
        elif i > 0:
            # Middle sections: remove first 4 (in prev transition) and last 4 (in next transition)
            if len(parsed_sections[i]) > 8:
                kept_segments = parsed_sections[i][4:-4]
                smoothed_sections.append('\n'.join([seg[1] for seg in kept_segments]))
            # else: entire section covered by transitions

    # Combine all parts
    final_script = '\n'.join([s for s in smoothed_sections if s.strip()])

    if show_progress:
        print(f"    ✓ All transitions smoothed (Total: {total_usage['input']} in, {total_usage['output']} out)")

    return final_script, total_usage


def run_multi_call_generation(topic, duration, word_count, research_context, source_documents,
                              web_source_count, style_template, style_description, language,
                              api_key, config, project_name, provider='elevenlabs'):
    """
    Main orchestrator for multi-call script generation.
    Returns final script text.

    Args:
        provider: 'elevenlabs' or 'cartesia' - determines emotion tag instructions
    """
    print("\n" + "="*60)
    print("MULTI-CALL SCRIPT GENERATION")
    print("="*60)

    total_usage = {'input': 0, 'output': 0}

    # Phase 1: Research
    research_summary = ""
    if web_source_count > 0:
        print("\n[PHASE 1/5] Web Research")
        research_summary = research_web_sources(
            topic=topic,
            research_context=research_context,
            source_count=web_source_count,
            api_key=api_key,
            config=config
        )
    else:
        print("\n[PHASE 1/5] Web Research - Skipped (0 sources requested)")

    # Phase 2: Document Processing
    doc_summary = ""
    if source_documents:
        print("\n[PHASE 2/5] Document Processing")
        doc_summary = process_documents_batched(
            documents_text=source_documents,
            project_name=project_name,
            api_key=api_key,
            config=config
        )
    else:
        print("\n[PHASE 2/5] Document Processing - Skipped (no documents)")

    # Phase 3: Outline Generation
    print("\n[PHASE 3/5] Outline Generation")
    outline, outline_usage = generate_outline(
        topic=topic,
        duration=duration,
        word_count=word_count,
        research_summary=research_summary,
        doc_summary=doc_summary,
        style_description=style_description,
        language=language,
        api_key=api_key,
        config=config
    )

    if not outline:
        print("✗ Failed to generate outline")
        return None

    if outline_usage:
        total_usage['input'] += outline_usage.input_tokens
        total_usage['output'] += outline_usage.output_tokens

    # Save outline for debugging
    outline_path = Path(f"./projects/{project_name}/debug/outline.txt")
    outline_path.parent.mkdir(parents=True, exist_ok=True)
    with open(outline_path, 'w', encoding='utf-8') as f:
        f.write(outline)
    print(f"    [DEBUG] Outline saved to: {outline_path}")

    # Phase 4: Script Generation
    print("\n[PHASE 4/5] Script Generation")
    sections, script_usage = generate_script_multi_call(
        topic=topic,
        duration=duration,
        word_count=word_count,
        outline=outline,
        style_template=style_template,
        language=language,
        api_key=api_key,
        config=config,
        provider=provider
    )

    if not sections:
        print("✗ Failed to generate script sections")
        return None

    if script_usage:
        total_usage['input'] += script_usage['input']
        total_usage['output'] += script_usage['output']

    # Phase 5: Lightweight Synthesis (smooth transitions between sections)
    if len(sections) > 1:
        print(f"\n[PHASE 5/5] Lightweight Synthesis ({len(sections)-1} transition(s))")
        final_script, synth_usage = synthesize_transitions(
            sections=sections,
            language=language,
            api_key=api_key,
            config=config
        )
        if synth_usage:
            total_usage['input'] += synth_usage['input']
            total_usage['output'] += synth_usage['output']
    else:
        print("\n[PHASE 5/5] Synthesis - Skipped (single section)")
        final_script = sections[0]

    # Post-process to ensure clean format
    print("\n[POST-PROCESS] Cleaning script format...")
    final_script = clean_script_format(final_script)
    print("    ✓ Removed markdown/titles/blank lines")

    print("\n" + "="*60)
    print("✓ MULTI-CALL GENERATION COMPLETE")
    print("="*60)
    word_count_actual = len(final_script.split())
    print(f"  Final script: {word_count_actual} words")
    print(f"  Total tokens - Input: {total_usage['input']}, Output: {total_usage['output']}")
    print("="*60)

    return final_script


def fetch_and_save_sources_separately(project_name, topic, api_key):
    """Make a separate Claude call to get research sources"""
    print("\n" + "="*60)
    print("FETCHING RESEARCH SOURCES")
    print("="*60)
    
    sources_prompt = f"""List all the research sources you used when creating the podcast about "{topic}".

Format each source as:
Title - URL

Examples:
Nature Study on Quantum Computing 2024 - https://nature.com/articles/...
MIT Technology Review: AI Advances - https://technologyreview.com/...

Only list sources, no commentary or explanations.
If you did not use external sources, respond with: "No external sources used."
"""
    
    try:
        response = anthropic.Anthropic(api_key=api_key).messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=2000,
            messages=[{"role": "user", "content": sources_prompt}]
        )
        
        sources_content = response.content[0].text.strip()
        
        # Save to sources folder (not scripts folder!)
        sources_file = Path(f"./projects/{project_name}/sources/{project_name}_sources.txt")
        sources_file.parent.mkdir(parents=True, exist_ok=True)
        
        with open(sources_file, 'w', encoding='utf-8') as f:
            f.write(f"Research Sources for {project_name}\n")
            f.write("="*60 + "\n\n")
            f.write(sources_content)
            f.write("\n\n")
            f.write("="*60 + "\n")
            f.write(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}\n")
        
        print(f"✓ Sources saved to: {sources_file}")
        return True
        
    except Exception as e:
        print(f"⚠ Could not fetch sources: {e}")
        return False


def extract_and_save_sources(script, project_name):
    """Legacy function - now just cleans up any sources Claude included.
    Sources are fetched separately via fetch_and_save_sources_separately()"""
    
    # Remove any sources that Claude included despite instructions
    sources_pattern = r'(?:^|\n)(?:\*\*)?(?:SOURCES FOUND:|Sources?:)(?:\*\*)?(.*?)(?:\n\n---|$)'
    script_clean = re.sub(sources_pattern, '', script, flags=re.DOTALL | re.IGNORECASE)
    
    # Also remove any remaining source lists at end
    script_clean = re.sub(r'\n\d+\.\s+.*?https?://.*', '', script_clean, flags=re.MULTILINE)
    
    return script_clean.strip()


def save_script(script, project_name, language_code, provider_tag, draft_number):
    """Save script with versioned filename including language, provider tag and timestamp"""
    timestamp = datetime.now().strftime('%Y-%m-%d_%H-%M')
    lang_upper = language_code.upper()
    filename = f"{project_name}_{lang_upper}_{timestamp}_{provider_tag}_draft{draft_number}.txt"
    path = Path(f"./projects/{project_name}/scripts/{filename}")
    with open(path, 'w', encoding='utf-8') as f:
        f.write(script)
    return path


def save_prompt(prompt, project_name, filename):
    """Save prompt to project folder"""
    path = Path(f"./projects/{project_name}/prompts/{filename}")
    with open(path, 'w', encoding='utf-8') as f:
        f.write(prompt)
    return path


def clean_script_for_audio(script):
    """Remove non-dialogue content before audio generation
    
    Removes:
    - Claude's meta-commentary preamble
    - Search quality checks
    - Search tags
    - Sources section
    - Markdown formatting
    - Stage directions
    """
    print("\n[INFO] Cleaning script for audio generation...")
    
    original_length = len(script)
    
    # CRITICAL: Remove Claude's meta-commentary at the start
    # Everything before the first "Speaker A:" or "Speaker B:"
    speaker_pattern = r'(?:Speaker [AB]:)'
    match = re.search(speaker_pattern, script)
    if match:
        # Found first speaker label - keep everything from there
        script = script[match.start():]
        print(f"[INFO] Removed Claude's preamble ({match.start()} chars)")
    
    # Remove search quality checks and search tags
    script = re.sub(r'<search_quality_check>.*?</search_quality_check>', '', script, flags=re.DOTALL)
    script = re.sub(r'<search_quality_score>.*?</search_quality_score>', '', script, flags=re.DOTALL)
    script = re.sub(r'<search>.*?</search>', '', script, flags=re.DOTALL)
    
    # Remove "I'll conduct research..." type preambles
    script = re.sub(r"(?:^|\n)I'?ll? (?:conduct|create|generate|search).*?(?=Speaker [AB]:|$)", '', script, flags=re.DOTALL|re.IGNORECASE)
    script = re.sub(r"(?:^|\n)Let me (?:conduct|create|generate|search).*?(?=Speaker [AB]:|$)", '', script, flags=re.DOTALL|re.IGNORECASE)
    script = re.sub(r"(?:^|\n)Now I'?ll? (?:conduct|create|generate).*?(?=Speaker [AB]:|$)", '', script, flags=re.DOTALL|re.IGNORECASE)
    
    # CRITICAL: Remove sources section FIRST - CUT EVERYTHING after "SOURCES FOUND:"
    # This must happen BEFORE removing "---" because there's often a "---" before sources
    print("[DEBUG] Checking for sources section...")
    
    sources_removed = False
    for pattern in [r'\n\s*SOURCES FOUND:', r'\n\s*\*\*SOURCES FOUND:\*\*', r'\n\s*##\s*SOURCES FOUND:']:
        match = re.search(pattern, script, re.IGNORECASE)
        if match:
            before_length = len(script)
            # Cut everything from the match position onwards
            script = script[:match.start()]
            after_length = len(script)
            print(f"[INFO] ✓ CUT SOURCES: Removed {before_length - after_length} chars after 'SOURCES FOUND:'")
            sources_removed = True
            break
    
    if not sources_removed:
        print("[WARNING] No 'SOURCES FOUND:' marker detected - check if script has sources section")
        # Try to find if there are numbered sources like "1. **Source**"
        if re.search(r'\n\d+\.\s+\*\*.*?\*\*', script):
            print("[WARNING] Found numbered sources but no 'SOURCES FOUND:' marker - may include sources in audio!")
    else:
        # Double-check sources are gone
        if re.search(r'\n\d+\.\s+\*\*.*?\*\*', script):
            print("[ERROR] Sources still present after removal! Check script format.")
        else:
            print("[INFO] ✓ Verified: No sources in cleaned script")
    
    # NOW remove "---" separators (often appear before sources section)
    script = re.sub(r'^-{3,}$', '', script, flags=re.MULTILINE)
    print("[DEBUG] Removed separator lines (---)")
    
    # Remove markdown headers
    script = re.sub(r'^#+\s+.*$', '', script, flags=re.MULTILINE)
    
    # Remove stage directions (but NOT audio tags!)
    script = re.sub(r'^\*[^\[]*\*$', '', script, flags=re.MULTILINE)
    
    # Remove word counts (various formats including markdown bold)
    script = re.sub(r'^\s*\*?\*?Word count:?\s*\d+\s*words?\*?\*?\s*$', '', script, flags=re.MULTILINE|re.IGNORECASE)
    script = re.sub(r'^\s*\*?\*?(?:Total|Approximate)?\s*(?:script\s+)?(?:length|count)?:?\s*~?\d+\s*words?\*?\*?\s*$', '', script, flags=re.MULTILINE|re.IGNORECASE)
    script = re.sub(r'Total script length:.*$', '', script, flags=re.MULTILINE|re.IGNORECASE)
    
    # Clean up extra blank lines
    script = re.sub(r'\n{3,}', '\n\n', script)
    script = script.strip()
    
    cleaned_length = len(script)
    removed = original_length - cleaned_length
    
    if removed > 0:
        print(f"[INFO] Removed {removed} characters of non-dialogue content")
    
    return script


def validate_template_quality(script):
    """Check if script uses dynamic features"""
    warnings = []
    
    if '[interrupting]' not in script.lower() and '[overlapping]' not in script.lower():
        warnings.append("⚠ No interruptions found - dialogue may sound too formal")
    
    emotion_tags = ['[excited]', '[curious]', '[skeptical]', '[surprised]', '[thoughtful]']
    if not any(tag.lower() in script.lower() for tag in emotion_tags):
        warnings.append("⚠ No emotional tags found - dialogue may lack energy")
    
    reaction_tags = ['[laughs]', '[chuckles]', '[sighs]', '[gasps]']
    if not any(tag.lower() in script.lower() for tag in reaction_tags):
        warnings.append("⚠ No reaction tags found - may sound robotic")
    
    if re.search(r'\bSie\b', script):
        warnings.append("⚠ Found 'Sie' form - should use informal 'Du' for friendly tone")
    
    if warnings:
        print("\n" + "="*60)
        print("SCRIPT QUALITY WARNINGS")
        print("="*60)
        for warning in warnings:
            print(warning)
        print("")
        print("Consider revising for better audio quality.")
        print("="*60)
        
        proceed = input("\nProceed anyway? (Y/n): ")
        if proceed.lower() == 'n':
            return False
    
    return True


def save_debug_chunk(chunk, chunk_num, project_name):
    """Save chunk content for debugging"""
    debug_dir = Path(f"./projects/{project_name}/debug")
    debug_file = debug_dir / f"chunk_{chunk_num}_content.json"
    
    with open(debug_file, 'w', encoding='utf-8') as f:
        json.dump(chunk, f, indent=2, ensure_ascii=False)
    
    log_debug(f"Chunk {chunk_num} saved to: {debug_file}")
    return debug_file


def parse_script_to_dialogue(script, voice_ids):
    """Parse script with Speaker A/B labels into ElevenLabs dialogue format
    
    CRITICAL: Preserves [audio tags] in square brackets for ElevenLabs v3
    """
    print("\n[DEBUG] Parsing script into dialogue format...")
    print(f"[DEBUG] Script length: {len(script)} characters")
    
    lines = script.split('\n')
    dialogue = []
    current_speaker = None
    current_text = []
    
    print("[DEBUG] First 10 lines of script:")
    for i, line in enumerate(lines[:10]):
        print(f"  {i}: {line[:80]}")
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
        
        is_speaker_a = any(marker in line.lower() for marker in 
                          ['speaker a:', '**speaker a', 'speaker a -'])
        is_speaker_b = any(marker in line.lower() for marker in 
                          ['speaker b:', '**speaker b', 'speaker b -'])
        
        if is_speaker_a:
            if current_text and current_speaker:
                dialogue.append({
                    'voice_id': voice_ids['speaker_a' if current_speaker == 'speaker_a' else 'speaker_b'],
                    'text': ' '.join(current_text).strip()
                })
                print(f"[DEBUG] Added {current_speaker} segment: {len(' '.join(current_text))} chars")
            
            current_speaker = 'speaker_a'
            # CRITICAL FIX: Don't strip [square brackets] - they're audio tags!
            text = line.split(':', 1)[-1].strip().replace('**', '').strip()
            current_text = [text] if text else []
            
        elif is_speaker_b:
            if current_text and current_speaker:
                dialogue.append({
                    'voice_id': voice_ids['speaker_a' if current_speaker == 'speaker_a' else 'speaker_b'],
                    'text': ' '.join(current_text).strip()
                })
                print(f"[DEBUG] Added {current_speaker} segment: {len(' '.join(current_text))} chars")
            
            current_speaker = 'speaker_b'
            # CRITICAL FIX: Don't strip [square brackets] - they're audio tags!
            text = line.split(':', 1)[-1].strip().replace('**', '').strip()
            current_text = [text] if text else []
            
        elif current_speaker:
            if not line.startswith('#') and not line.startswith('---'):
                current_text.append(line)
    
    if current_text and current_speaker:
        voice_id = voice_ids['speaker_a'] if current_speaker == 'speaker_a' else voice_ids['speaker_b']
        dialogue.append({'voice_id': voice_id, 'text': ' '.join(current_text).strip()})
        print(f"[DEBUG] Added final {current_speaker} segment: {len(' '.join(current_text))} chars")
    
    print(f"[DEBUG] Total dialogue segments: {len(dialogue)}")
    
    if not dialogue:
        print("[ERROR] No dialogue segments found!")
        print("[ERROR] Script may not have proper Speaker A: / Speaker B: labels")
        print("\n[HELP] Script should look like:")
        print("Speaker A: Hello, welcome to the podcast!")
        print("Speaker B: Thanks for having me!")
        return None
    
    return dialogue


def chunk_dialogue(inputs, max_chars=4500):
    """Split dialogue inputs into chunks under character limit"""
    chunks = []
    current_chunk = []
    current_length = 0
    
    for item in inputs:
        item_length = len(item['text'])
        
        if current_length + item_length > max_chars and current_chunk:
            chunks.append(current_chunk)
            current_chunk = [item]
            current_length = item_length
        else:
            current_chunk.append(item)
            current_length += item_length
    
    if current_chunk:
        chunks.append(current_chunk)
    
    return chunks


def generate_audio(script, config, language_code, provider_name, mode='prototype', speed=1.0, project_name=None):
    """Generate audio using specified TTS provider"""
    
    # Get provider instance
    provider = get_provider_instance(provider_name, config)
    if not provider:
        return None, 0
    
    # Get language mapping
    language_map = {'de': 'german', 'en': 'english', 'nl': 'dutch'}
    language = language_map.get(language_code, 'english')
    
    # Get voice IDs from provider config
    provider_config = config['providers'][provider_name]
    if language not in provider_config.get('voices', {}):
        print(f"ERROR: No voices for {language} in provider {provider_name}")
        return None, 0
    
    voice_config = provider_config['voices'][language]
    voice_ids = {
        'speaker_a': voice_config.get('speaker_a_female') or voice_config.get('speaker_a_male'),
        'speaker_b': voice_config.get('speaker_b_male') or voice_config.get('speaker_b_female')
    }
    
    print(f"[DEBUG] Voice A: {voice_ids['speaker_a']}")
    print(f"[DEBUG] Voice B: {voice_ids['speaker_b']}")
    
    # Generate using provider
    return provider.generate_audio(script, voice_ids, mode, speed, project_name)


# Original generate_audio function preserved for backward compatibility
def generate_audio_legacy(script, config, language_code, mode='prototype', speed=1.0, project_name=None):
    """Call ElevenLabs Text-to-Dialogue API with enhanced error logging"""
    print(f"\nGenerating audio in {mode.upper()} mode...")
    
    api_key = os.getenv('ELEVENLABS_API_KEY')
    if not api_key:
        print("ERROR: ELEVENLABS_API_KEY not found in environment")
        return None
    
    # Map language codes to config keys for ALL supported languages
    language_map = {
        'de': 'german',
        'en': 'english',
        'nl': 'dutch'
    }
    language = language_map.get(language_code, 'english')
    
    voice_ids = {
        'speaker_a': config['languages'][language]['elevenlabs_voices']['speaker_a_female'],
        'speaker_b': config['languages'][language]['elevenlabs_voices']['speaker_b_male']
    }
    
    print(f"[DEBUG] Language: {language.upper()} ({language_code})")
    print(f"[DEBUG] Using voices: Speaker A = {voice_ids['speaker_a']}, Speaker B = {voice_ids['speaker_b']}")
    
    dialogue = parse_script_to_dialogue(script, voice_ids)
    
    if not dialogue:
        print("\n" + "="*60)
        print("SCRIPT FORMAT ERROR")
        print("="*60)
        print("The script doesn't have proper Speaker A: / Speaker B: labels.")
        print("")
        print("Options to fix:")
        print("1. Go back and ask Claude to revise with proper format")
        print("2. Manually edit the script file to add labels")
        print("3. Use 'Edit script and regenerate' option")
        print("="*60)
        return None
    
    inputs = [{
        "text": seg['text'], 
        "voice_id": seg['voice_id'],
        "voice_settings": {"speed": speed}
    } for seg in dialogue]
    
    total_length = sum(len(item['text']) for item in inputs)
    print(f"[DEBUG] Total dialogue length: {total_length} characters")
    
    if total_length > 5000:
        print(f"[INFO] Content exceeds 5000 character limit, splitting into chunks...")
        chunks = chunk_dialogue(inputs, max_chars=4500)
        print(f"[INFO] Split into {len(chunks)} chunks")
    else:
        chunks = [inputs]
        print(f"[INFO] Content fits in single request")
    
    url = "https://api.elevenlabs.io/v1/text-to-dialogue"
    headers = {"xi-api-key": api_key, "Content-Type": "application/json"}
    
    audio_parts = []
    
    for i, chunk in enumerate(chunks, 1):
        chunk_length = sum(len(item['text']) for item in chunk)
        print(f"\n[DEBUG] Chunk {i}/{len(chunks)}: {len(chunk)} segments, {chunk_length} characters")
        
        # Save chunk for debugging
        if project_name:
            debug_file = save_debug_chunk(chunk, i, project_name)
            print(f"[DEBUG] Chunk {i} saved to: {debug_file}")
        
        payload = {"inputs": chunk}
        
        # Retry logic with detailed error messages
        max_retries = 3
        retry_delay = 2
        
        for attempt in range(max_retries):
            try:
                log_debug(f"Attempt {attempt + 1}/{max_retries} for chunk {i}/{len(chunks)}")
                
                if attempt > 0:
                    print(f"[RETRY] Attempt {attempt + 1}/{max_retries} for chunk {i}/{len(chunks)}...")
                    time.sleep(retry_delay * attempt)
                
                print(f"Sending chunk {i}/{len(chunks)} to ElevenLabs...")
                
                # Make the request with timeout
                log_debug(f"POST {url}")
                log_debug(f"Payload size: {len(json.dumps(payload))} bytes")
                
                response = requests.post(url, headers=headers, json=payload, stream=True, timeout=120)
                
                log_debug(f"Response status: {response.status_code}")
                
                if response.status_code != 200:
                    error_body = response.text
                    print(f"\n[ERROR] Status {response.status_code}: {error_body}")
                    
                    # If 500 error, retry
                    if response.status_code == 500 and attempt < max_retries - 1:
                        print(f"[INFO] Server error on attempt {attempt + 1}, retrying in {retry_delay * (attempt + 1)} seconds...")
                        continue
                    
                    # If other error or last attempt, raise
                    response.raise_for_status()
                
                # Success - collect audio
                chunk_audio = b''
                bytes_received = 0
                for data in response.iter_content(chunk_size=8192):
                    if data:
                        chunk_audio += data
                        bytes_received += len(data)
                
                log_debug(f"Received {bytes_received} bytes")
                
                audio_parts.append(chunk_audio)
                print(f"✓ Chunk {i}/{len(chunks)} generated ({len(chunk_audio) / 1024 / 1024:.1f} MB)")
                break  # Success, exit retry loop
                
            except requests.exceptions.Timeout as e:
                print(f"\n[ERROR] Timeout after 120 seconds on chunk {i}")
                if attempt < max_retries - 1:
                    print(f"[INFO] Retrying in {retry_delay * (attempt + 1)} seconds...")
                    continue
                else:
                    print(f"\n✗ Failed after {max_retries} attempts: Timeout")
                    return None
                    
            except requests.exceptions.RequestException as e:
                print(f"\n[ERROR] Request exception on chunk {i}: {type(e).__name__}")
                print(f"[ERROR] Details: {str(e)}")
                
                if hasattr(e, 'response') and e.response is not None:
                    print(f"[ERROR] Response body: {e.response.text}")
                
                if attempt < max_retries - 1:
                    print(f"[INFO] Retrying in {retry_delay * (attempt + 1)} seconds...")
                    continue
                else:
                    print(f"\n✗ Failed after {max_retries} attempts")
                    print(f"\n[DEBUG] Chunk {i} content saved to:")
                    print(f"  projects/{project_name}/debug/chunk_{i}_content.json")
                    return None
            
            except Exception as e:
                print(f"\n[ERROR] Unexpected exception on chunk {i}: {type(e).__name__}")
                print(f"[ERROR] Details: {str(e)}")
                
                if attempt < max_retries - 1:
                    print(f"[INFO] Retrying in {retry_delay * (attempt + 1)} seconds...")
                    continue
                else:
                    return None
    
    # Concatenate all audio chunks
    if len(audio_parts) > 1:
        print(f"\n[INFO] Concatenating {len(audio_parts)} audio chunks...")
        audio_data = b''.join(audio_parts)
    else:
        audio_data = audio_parts[0]
    
    print(f"✓ Complete audio generated ({len(audio_data) / 1024 / 1024:.1f} MB)")
    print(f"[USAGE] ElevenLabs - {total_length} characters processed")
    
    return audio_data, total_length


def read_text_file(filepath):
    """Read plain text file with verbose feedback"""
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            content = f.read()
        
        lines = content.split('\n')
        print(f"      [Text file: {len(lines)} lines, {len(content)} chars]")
        return content
    except Exception as e:
        return f"[Error reading {filepath}: {str(e)}]"


def read_docx_file(filepath):
    """Read DOCX file with verbose feedback"""
    if not DOCX_AVAILABLE:
        return "[python-docx not installed - run: pip install python-docx]"
    
    try:
        doc = DocxDocument(filepath)
        num_paragraphs = len(doc.paragraphs)
        print(f"      [DOCX: {num_paragraphs} paragraphs detected]")
        
        text = []
        for para in doc.paragraphs:
            if para.text.strip():
                text.append(para.text)
        
        print(f"      [Extracted {len(text)} non-empty paragraphs]")
        return '\n'.join(text)
    except Exception as e:
        return f"[Error reading {filepath}: {str(e)}]"


def read_pdf_file(filepath):
    """Read PDF file with verbose feedback"""
    if not PDF_AVAILABLE:
        return "[PyPDF2 not installed - run: pip install PyPDF2]"
    
    try:
        text = []
        with open(filepath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            num_pages = len(reader.pages)
            print(f"      [PDF: {num_pages} pages detected]")
            
            for page_num, page in enumerate(reader.pages, 1):
                page_text = page.extract_text()
                if page_text.strip():
                    text.append(page_text)
                    print(f"      [Page {page_num}/{num_pages}: {len(page_text)} chars]", end='\r')
            
            print()  # New line after progress
        return '\n'.join(text)
    except Exception as e:
        return f"[Error reading {filepath}: {str(e)}]"


def read_pptx_file(filepath):
    """Read PPTX file with verbose feedback"""
    if not PPTX_AVAILABLE:
        return "[python-pptx not installed - run: pip install python-pptx]"
    
    try:
        prs = Presentation(filepath)
        num_slides = len(prs.slides)
        print(f"      [PPTX: {num_slides} slides detected]")
        
        text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            text.append(f"[Slide {slide_num}]")
            shape_count = 0
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
                    shape_count += 1
            print(f"      [Slide {slide_num}/{num_slides}: {shape_count} text elements]", end='\r')
        
        print()  # New line after progress
        return '\n'.join(text)
    except Exception as e:
        return f"[Error reading {filepath}: {str(e)}]"


def read_source_document(filepath):
    """Read document based on file extension"""
    path = Path(filepath)
    ext = path.suffix.lower()
    
    if ext in ['.txt', '.md']:
        return read_text_file(filepath)
    elif ext == '.docx':
        return read_docx_file(filepath)
    elif ext == '.pdf':
        return read_pdf_file(filepath)
    elif ext == '.pptx':
        return read_pptx_file(filepath)
    else:
        return f"[Unsupported file type: {ext}]"


def list_source_files(project_name):
    """List available source files in project sources folder"""
    sources_path = Path(f"./projects/{project_name}/sources")
    if not sources_path.exists():
        return []
    
    supported_exts = ['.txt', '.md', '.docx', '.pdf', '.pptx']
    files = []
    for ext in supported_exts:
        files.extend(sources_path.glob(f"*{ext}"))
    
    return sorted([f for f in files if f.name not in ['research_context.txt', 'sources_list.txt']])


def process_source_documents(project_name):
    """Check and process source documents before script generation"""
    sources_path = Path(f"./projects/{project_name}/sources")
    sources_path.mkdir(parents=True, exist_ok=True)
    
    while True:
        print("\n" + "="*60)
        print("SOURCE DOCUMENTS CHECK")
        print("="*60)
        
        files = list_source_files(project_name)
        
        if files:
            print(f"Found {len(files)} document(s) in sources folder:")
            for f in files:
                print(f"  - {f.name}")
        else:
            print("No source documents found in sources folder.")
        
        print("\nOptions:")
        print("  1. Proceed (use existing documents if any)")
        print("  2. List current documents")
        print("  3. Add new source files")
        
        choice = input("\nChoice (1-3, default=1): ").strip() or "1"
        
        if choice == "1":
            # Read all documents and return combined text
            if not files:
                print("\n[INFO] No source documents - proceeding with web research only")
                return ""
            
            print("\n[INFO] Reading source documents...")
            combined_text = []
            success_count = 0
            
            for file in files:
                print(f"  Reading: {file.name}...")
                content = read_source_document(file)
                if not content.startswith("[Error") and not content.startswith("["):
                    combined_text.append(f"\n\n### SOURCE: {file.name}\n\n{content}")
                    success_count += 1
                    print(f"    ✓ Successfully read ({len(content)} chars)")
                else:
                    print(f"    ✗ {content}")
            
            print(f"\n[INFO] Successfully read {success_count}/{len(files)} documents")
            return '\n'.join(combined_text) if combined_text else ""
        
        elif choice == "2":
            # List documents (filenames only, no content)
            if not files:
                print("\n[INFO] No documents to list")
                continue
            
            print("\n" + "="*60)
            print("DOCUMENTS IN SOURCES FOLDER")
            print("="*60)
            print(f"\nLocation: {sources_path.absolute()}\n")
            
            for i, file in enumerate(files, 1):
                file_size = file.stat().st_size
                size_kb = file_size / 1024
                print(f"  {i}. {file.name} ({size_kb:.1f} KB)")
            
            print(f"\nTotal: {len(files)} document(s)")
            input("\nPress Enter to continue...")
        
        elif choice == "3":
            # Provide instructions for adding files
            print("\n" + "="*60)
            print("ADD SOURCE DOCUMENTS")
            print("="*60)
            print(f"\nLocation: {sources_path.absolute()}")
            print("\nSupported formats:")
            print("  - Text: .txt, .md")
            print("  - Documents: .docx")
            print("  - PDF: .pdf")
            print("  - Presentations: .pptx")
            print("\nInstructions:")
            print("  1. Copy your files to the sources folder above")
            print("  2. Press Enter when done")
            print("\nNote: Files named 'research_context.txt' and 'sources_list.txt'")
            print("      are reserved and will be ignored.")
            
            input("\nPress Enter when files are ready...")
        
        else:
            print("Invalid choice")


def save_audio(audio_data, project_name, topic, language_code, provider_tag, mode, speed, config, is_test_mode=False, topic_tag=None):
    """Save audio file with project name, topic, language, and provider tag"""
    date = datetime.now().strftime('%Y-%m-%d')
    safe_topic = topic.replace('/', '-').replace('\\', '-')
    if is_test_mode and topic_tag:
        # test_de_2025-11-29_road-prsd_CRTS_OS1.05_MS1.00_FS1.10_PROTOTYPE.mp3
        # Add speed tags for test mode: OS (overall), MS (male), FS (female)
        speed_tag = f"OS{speed:.2f}"
        
        # Get male/female speeds from config if available
        male_speed = speed
        female_speed = speed
        if config and 'speed_adjustments' in config and isinstance(config.get('speed_adjustments'), dict):
            male_speed = speed * config['speed_adjustments'].get('speaker_b_male', 1.0)
            female_speed = speed * config['speed_adjustments'].get('speaker_a_female', 1.0)
        
        ms_tag = f"MS{male_speed:.2f}"
        fs_tag = f"FS{female_speed:.2f}"
        
        filename = f"{project_name.lower()}_{language_code}_{date}_{topic_tag}_{provider_tag}_{speed_tag}_{ms_tag}_{fs_tag}_{mode.upper()}.mp3"
    else:
        # Normal filename
        safe_topic = topic.replace('/', '-').replace('\\', '-')
        filename = f"{project_name}_{safe_topic}_{language_code}_{date}_{provider_tag}_{mode.upper()}.mp3"
    path = Path(f"./projects/{project_name}/audio/{filename}")
    
    # Ensure directory exists
    path.parent.mkdir(parents=True, exist_ok=True)
    
    with open(path, 'wb') as f:
        f.write(audio_data)
    
    print(f"[DEBUG] Saved {len(audio_data)} bytes to {path}")
    return path




def extract_provider_from_filename(filename):
    """Extract provider tag from script filename"""
    if '_11LB_' in filename:
        return 'elevenlabs', '11LB'
    elif '_CRTS_' in filename:
        return 'cartesia', 'CRTS'
    else:
        return 'elevenlabs', '11LB'  # Default


def get_provider_instance(provider_name, config):
    """Get TTS provider instance"""
    if provider_name not in config.get('providers', {}):
        print(f"ERROR: Provider '{provider_name}' not in config")
        return None
    
    provider_config = config['providers'][provider_name]
    api_key_env = provider_config.get('api_key_env')
    api_key = os.getenv(api_key_env)
    
    if not api_key:
        print(f"ERROR: {api_key_env} not found in config/.env")
        return None
    
    if provider_name == 'elevenlabs':
        return ElevenLabsProvider(api_key, provider_config)
    elif provider_name == 'cartesia':
        return CartesiaProvider(api_key, provider_config)
    else:
        print(f"ERROR: Unknown provider '{provider_name}'")
        return None


def inject_provider_instructions(template_content, provider_instance):
    """Inject provider-specific instructions into template"""
    provider_instructions = provider_instance.get_template_instructions()
    
    # Inject after the CRITICAL section
    marker = "===================================\nAVAILABLE AUDIO TAGS"
    if marker in template_content:
        template_content = template_content.replace(
            marker,
            provider_instructions + "\n" + marker
        )
    
    return template_content

def main():
    """Main pipeline orchestration"""
    print("=== AI Podcast Pipeline v3.0 (Enhanced Debug) ===\n")
    
    if DEBUG_VERBOSE:
        print("[VERBOSE MODE ENABLED - Detailed logging active]\n")
    
    config = load_config()
    
    anthropic_key = os.getenv('ANTHROPIC_API_KEY')
    if not anthropic_key:
        print("ERROR: ANTHROPIC_API_KEY not found in config/.env")
        return
    
    # 1. Project setup
    project_name = get_user_input("Enter project name (alphanumeric, for folders/filenames)").replace(' ', '_')
    
    # Detect test mode
    is_test_mode = (project_name.lower() == 'test')
    
    if is_test_mode:
        print("\n" + "="*60)
        print("TEST MODE ACTIVATED")
        print("="*60)
        print("Voice tuning test scenarios (~1-1.5 min, max 500 words)")
        print("="*60)
    # Test mode: scenario selection BEFORE normal topic
    if is_test_mode:
        # Scenario type selection
        print("\nSelect test scenario type:")
        scenario_type = get_user_input("", [
            "Road trip argument",
            "Cooking disaster",
            "Movie scene analysis",
            "Random scenario (Claude chooses interesting situation)"
        ])
        
        scenario_map = ['road', 'cook', 'mvie', 'rndm']
        selected_scenario = scenario_map[scenario_type]
        
        # Scenario detail selection
        if selected_scenario == 'road':
            print("\nRoad trip argument scenarios:")
            topic_choice = get_user_input("", [
                "Istanbul heavy traffic (Bosphorus bridge chaos)",
                "Paris late night (Arc de Triomphe roundabout)",
                "UK rural countryside (sheep blocking narrow roads)",
                "German Autobahn (no speed limit construction zone)",
                "Tokyo side streets (wrong GPS coordinates)",
                "Random (Claude searches and picks most interesting)"
            ])
            
            topics = {
                0: ('Istanbul Bosphorus bridge traffic chaos rush hour', 'istb'),
                1: ('Paris late night Arc de Triomphe roundabout confusion', 'prsd'),
                2: ('UK countryside single track roads with sheep blocking', 'ukrl'),
                3: ('German Autobahn no speed limit construction zone chaos', 'autb'),
                4: ('Tokyo narrow side streets GPS coordinates wrong', 'toky'),
                5: ('RANDOM', 'RANDOM')
            }
            
            topic_description, topic_abbrev = topics[topic_choice]
            
            if topic_abbrev == 'RANDOM':
                print("\n[INFO] Searching for interesting driving scenarios...")
                # Use Claude to search and select
                search_prompt = "Search for interesting unusual driving situations and traffic scenarios in cities worldwide. Return the most interesting one for a dialogue."
                # For now, use a default
                topic_description = "Venice narrow canals wrong boat GPS coordinates"
                topic_abbrev = "veni"
                print(f"[SELECTED] {topic_description}")
        
        elif selected_scenario == 'cook':
            print("\nCooking disaster scenarios:")
            topic_choice = get_user_input("", [
                "Pizza from scratch (dough disaster)",
                "Birthday cake (three-layer collapse)",
                "Christmas turkey (timing disaster dry meat)",
                "Chocolate soufflé (falling flat temperature fail)",
                "Sushi rice (first time sticky disaster)",
                "Random (Claude searches and picks most interesting)"
            ])
            
            topics = {
                0: ('Making authentic Neapolitan pizza from scratch dough disaster', 'pizz'),
                1: ('Three-layer birthday cake collapse frosting disaster', 'cake'),
                2: ('Christmas turkey timing disaster dry overcooked meat', 'xmas'),
                3: ('French chocolate soufflé falling flat temperature fail', 'souf'),
                4: ('First time making sushi rice sticky disaster', 'sush'),
                5: ('RANDOM', 'RANDOM')
            }
            
            topic_description, topic_abbrev = topics[topic_choice]
            
            if topic_abbrev == 'RANDOM':
                print("\n[INFO] Searching for cooking disasters...")
                topic_description = "Homemade pasta dough too sticky disaster"
                topic_abbrev = "past"
                print(f"[SELECTED] {topic_description}")
        
        elif selected_scenario == 'mvie':
            print("\nMovie scene analysis scenarios:")
            topic_choice = get_user_input("", [
                "Pulp Fiction (Royale with Cheese dialogue)",
                "Matrix (red pill blue pill philosophy)",
                "Star Trek (Borg resistance is futile)",
                "My Neighbor Totoro (cat bus magic scene)",
                "Inception (spinning top ending debate)",
                "Random (Claude searches and picks most interesting)"
            ])
            
            topics = {
                0: ('Pulp Fiction Royale with Cheese metric system dialogue', 'pulp'),
                1: ('Matrix red pill blue pill choice meaning philosophy', 'mtrx'),
                2: ('Star Trek Borg resistance is futile scene analysis', 'borg'),
                3: ('My Neighbor Totoro cat bus scene magic realism', 'totr'),
                4: ('Inception spinning top ending interpretation debate ambiguity', 'incr'),
                5: ('RANDOM', 'RANDOM')
            }
            
            topic_description, topic_abbrev = topics[topic_choice]
            
            if topic_abbrev == 'RANDOM':
                print("\n[INFO] Searching for famous movie scenes...")
                topic_description = "Casablanca play it again Sam misquote"
                topic_abbrev = "casa"
                print(f"[SELECTED] {topic_description}")
        
        else:  # rndm
            print("\n[INFO] Claude will search and select most interesting scenario...")
            # Let Claude pick everything
            topic_description = "Venice canals GPS navigation disaster"
            topic_abbrev = "veni"
            selected_scenario = 'road'  # Default to road for context
            print(f"[SELECTED] Road trip: {topic_description}")
        
        # Build topic_tag
        topic_tag = f"{selected_scenario}-{topic_abbrev}"
        topic = topic_description  # Use for display
        
        print(f"\nTopic set to: {topic_description}")
        print(f"File tag: {topic_tag}\n")
    
    else:
        # Normal mode
        topic = get_user_input("Enter podcast topic")
        topic_tag = None
        selected_scenario = None
        topic_description = None
    
    while True:
        try:
            duration = int(get_user_input("Target duration in minutes"))
            if duration > 0:
                break
            print("Duration must be positive")
        except ValueError:
            print("Please enter a valid number")
    
    word_count = duration * 222
    print(f"Calculated word count: ~{word_count} words (222 words/min)")

    # Web source count (for multi-call mode)
    gen_config = config.get('script_generation', {})
    if gen_config.get('enable_multi_call', False) and not is_test_mode:
        default_sources = gen_config.get('default_web_sources', 8)
        max_sources = gen_config.get('max_web_sources', 20)
        print(f"\nWeb research sources (more sources = more thorough, more API calls)")
        source_input = input(f"Number of sources to research (1-{max_sources}, default {default_sources}): ").strip()
        if source_input:
            try:
                web_source_count = int(source_input)
                web_source_count = max(1, min(max_sources, web_source_count))
            except ValueError:
                web_source_count = default_sources
        else:
            web_source_count = default_sources
        print(f"Will research {web_source_count} web sources")
    else:
        web_source_count = 0  # Test mode or single-call mode

    # 2-4. Style, language, mode selection
    styles = list(config['styles'].keys())
    style_names = [config['styles'][s]['description'] for s in styles]
    style_idx = get_user_input("\nSelect style", style_names)
    selected_style = styles[style_idx]
    
    languages = list(config['languages'].keys())
    language_names = [config['languages'][l]['name'] for l in languages]
    lang_idx = get_user_input("\nSelect language", language_names)
    selected_language = languages[lang_idx]
    language_code = config['languages'][selected_language]['code']

    # Recalculate word count using language default speed
    default_speed = config['languages'][selected_language].get('speed', 1.0)
    word_count = int(duration * 222 * default_speed)
    print(f"Adjusted word count: ~{word_count} words (for {default_speed} speed)")

    # 4b. TTS Provider Selection (BEFORE script generation for provider-optimized emotion tags)
    print("\n" + "="*60)
    print("TTS PROVIDER SELECTION")
    print("="*60)
    print("Select provider BEFORE script generation for optimized emotion tags.")
    provider_options = [
        "Cartesia (fast, affordable, 5 core emotions)",
        "ElevenLabs (premium, interruptions, overlapping)"
    ]
    provider_idx = get_user_input("\nSelect TTS provider", provider_options)
    selected_provider = "cartesia" if provider_idx == 0 else "elevenlabs"
    provider_tag = "CRTS" if selected_provider == "cartesia" else "11LB"
    print(f"\n[INFO] Selected: {selected_provider.upper()}")
    print("[INFO] Script will use provider-optimized emotion tags.")

    # 5. Create project structure
    print(f"\nCreating project folder: ./projects/{project_name}/")
    project_path = create_project_structure(project_name)
    print(f"  ✓ Created subdirectories")

    # 5b. CHECK FOR EXISTING SCRIPTS (any provider)
    if is_test_mode:
        # For tests: filter by scenario tag
        pattern = f"{project_name}_{language_code.upper()}_*_{topic_tag}_*_draft*.txt"
    else:
        # Normal pattern - any provider
        pattern = f"{project_name}_{language_code.upper()}_*_draft*.txt"

    existing_scripts = sorted(
        list(Path(f"./projects/{project_name}/scripts").glob(pattern)),
        key=lambda x: x.stat().st_mtime,
        reverse=True
    )
    
    script_ready = False

    if existing_scripts:
        print(f"\n⚠️  Found {len(existing_scripts)} existing script(s):")

        # Show all scripts
        for i, script_file in enumerate(existing_scripts, 1):
            print(f"    {i}. {script_file.name}")

        # Build options list
        options = [f"Use script #{i+1}" for i in range(len(existing_scripts))]
        options.append("Generate new script (continue with research/prompt setup)")
        options.append("Cancel")

        action = get_user_input("", options)

        # User selected an existing script (indices 0 to len-1)
        if action < len(existing_scripts):
            selected_script = existing_scripts[action]
            print(f"\n[INFO] Loading: {selected_script.name}")
            with open(selected_script, 'r', encoding='utf-8') as f:
                script = f.read()
            script_path = selected_script
            script_ready = True
            print("[INFO] Skipping to audio generation...\n")

        # User selected "Generate new script"
        elif action == len(existing_scripts):
            print("\n[INFO] Continuing with new script generation...\n")
            script_ready = False

        # User selected "Cancel"
        else:
            print("Cancelled")
            return

    # If using existing script, get TTS settings and skip to audio
    if script_ready:
        # TTS CONFIGURATION for existing script
        print("\n" + "="*60)
        print("AUDIO CONFIGURATION")
        print("="*60)
        provider_options = [
            "Cartesia (fast, affordable, 5 core emotions)",
            "ElevenLabs (premium, interruptions, overlapping)"
        ]
        provider_idx = get_user_input("\nSelect TTS provider", provider_options)
        selected_provider = "cartesia" if provider_idx == 0 else "elevenlabs"
        provider_tag = "CRTS" if selected_provider == "cartesia" else "11LB"

        print(f"\n[INFO] Selected: {selected_provider.upper()}")

        if selected_provider == 'elevenlabs':
            print("\n[INFO] ElevenLabs supports quality tiers:")
            print("  - Prototype: 64kbps (lower cost, testing)")
            print("  - Production: 128kbps+ (full quality)")
        elif selected_provider == 'cartesia':
            print("\n[INFO] Cartesia note: Always generates full quality")

        mode_idx = get_user_input("\nSelect mode", [
            "Prototype (lower quality, reduced cost for testing)",
            "Production (full quality)"
        ])
        mode = "prototype" if mode_idx == 0 else "production"

        default_speed = config['languages'][selected_language]['speed']
        speed_input = input(f"\nSpeech speed (0.7-1.2, default {default_speed}, Enter to use default): ").strip()
        if speed_input:
            try:
                speed = float(speed_input)
                speed = max(0.7, min(1.2, speed))
            except ValueError:
                speed = default_speed
        else:
            speed = default_speed
        print(f"Using speed: {speed}")
    
    # 5c. Research context (only if generating new)
    if not script_ready:
        
        if is_test_mode:
            # TEST MODE: Skip research context menu, use empty context
            research_context = ""
            print("\n[TEST MODE] Skipping research context (not needed for test scenarios)")
            
        else:
            # NORMAL MODE: Full research context flow
            research_context_file = project_path / "sources" / "research_context.txt"
            print(f"\n✓ Research context file: {research_context_file}")
            # Show what's being used
            default_template = Path("templates/research_contexts/default.txt")
            if default_template.exists() and not (project_path / "sources" / "research_context.txt").exists():
                print("  (Using default template from templates/research_contexts/default.txt)")
            elif (project_path / "sources" / "research_context.txt").exists():
                # Check if it's different from default (i.e., project-specific)
                with open(project_path / "sources" / "research_context.txt", 'r') as f:
                    current_content = f.read()
            
                is_customized = "{project_name}" not in current_content  # Simple check
                if is_customized:
                    print("  (Using project-specific research context)")
                else:
                    print("  (Using default template)")
        
            # Offer choices
            edit_choice = get_user_input("\nResearch context options", [
                "Use as-is (proceed with current context)",
                "Edit current context (customize for this project)",
                "Reset to default template (if you made mistakes)",
                "Show current context"
            ])
    
            if edit_choice == 1:
                print("\nOpening research context in your text editor...")
                subprocess.run([get_text_editor(), str(research_context_file)])
                print("✓ Research context updated (now project-specific)")
            elif edit_choice == 2:
                if default_template.exists():
                    print("\nResetting to default template...")
                    with open(default_template, 'r') as f:
                        template_content = f.read()
                    with open(research_context_file, 'w') as f:
                        f.write(template_content.replace("{project_name}", project_name))
                    print("✓ Reset to default template")
                else:
                    print("⚠ No default template found at templates/research_contexts/default.txt")
            elif edit_choice == 3:
                print("\n" + "="*60)
                with open(research_context_file, 'r') as f:
                    print(f.read())
                print("="*60)
                input("\nPress Enter to continue...")
    
            with open(research_context_file, 'r', encoding='utf-8') as f:
                research_context = f.read()
    
        # 6. Prompt handling
        if is_test_mode:
            # TEST MODE: Load test template directly, skip menu
            print("\n[TEST MODE] Loading test template...")
            
            # Map style to template filename prefix
            # selected_style is like "popular_scientific", template files are "popular_science"
            style_to_template = {
                'popular_scientific': 'popular_science',
                'technical_deep_dive': 'technical_deep_dive',
                'news_brief': 'news_brief'
            }
            style_key = style_to_template.get(selected_style, selected_style)
            
            template_file = f"templates/{style_key}_{selected_language}_TEST.txt"
            if not Path(template_file).exists():
                print(f"[WARNING] Test template not found: {template_file}")
                print(f"[INFO] Falling back to regular template")
                template_file = f"templates/{style_key}_{selected_language}_dynamic.txt"
            
            print(f"Using template: {template_file}")
            
            with open(template_file, 'r', encoding='utf-8') as f:
                template = f.read()
            
            # Inject scenario context
            scenario_context = build_scenario_context(selected_scenario, topic_description, selected_language)
            template = template.replace('{SCENARIO_DESCRIPTION}', topic_description)
            template = template.replace('{SCENARIO_SPECIFIC_INSTRUCTIONS}', scenario_context)
            
            prompt = template  # Test template is already complete
            
        else:
            # NORMAL MODE: Full prompt template menu
            variables = {
                'duration': duration,
                'word_count': word_count,
                'topic': topic,
                'project_name': project_name
            }
    
            prompt_choice = get_user_input("\nPrompt template options", [
                f"Use default template ({selected_style} / {selected_language})",
                "Load existing template from this project's prompts folder",
                "Copy template from templates folder to project and customize",
                "Edit the chosen template before generating",
                "Start with blank prompt"
            ])
    
            if prompt_choice == 0:
                # Use default template
                template_file = config['styles'][selected_style]['default_template_file']
                template_file = template_file.replace('{language}', selected_language)
                if Path(template_file).exists():
                    prompt = load_template(template_file, variables)
                else:
                    print(f"WARNING: Template {template_file} not found")
                    prompt = f"Create a {duration}-minute podcast script about '{topic}'."
            
            elif prompt_choice == 1:
                # Load existing from project
                project_prompts = list(Path(f"./projects/{project_name}/prompts/").glob("*.txt"))
                if project_prompts:
                    prompt_names = [p.name for p in project_prompts]
                    prompt_idx = get_user_input("Select prompt file", prompt_names)
                    prompt = load_template(project_prompts[prompt_idx], variables)
                else:
                    print("No saved prompts found in project")
                    prompt = f"Create a {duration}-minute podcast script about '{topic}'."
            
            elif prompt_choice == 2:
                # Copy global template to project
                templates = list(Path("./templates/").glob("*.txt"))
                if templates:
                    template_names = [t.name for t in templates]
                    template_idx = get_user_input("Select template to copy", template_names)
                    prompt = load_template(templates[template_idx], variables)
                    save_prompt(prompt, project_name, "copied_template.txt")
                    print(f"Template copied to project")
                else:
                    prompt = f"Create a {duration}-minute podcast script about '{topic}'."
            
            elif prompt_choice == 3:
                # Edit template before generating
                template_file = config['styles'][selected_style]['default_template_file']
                template_file = template_file.replace('{language}', selected_language)
                if Path(template_file).exists():
                    prompt = load_template(template_file, variables)
                else:
                    prompt = f"Create a {duration}-minute podcast script about '{topic}'."
        
                save_prompt(prompt, project_name, "edited_prompt.txt")
                prompt_file = Path(f"./projects/{project_name}/prompts/edited_prompt.txt")
                print(f"\nOpening {prompt_file} for editing...")
                subprocess.run([get_text_editor(), str(prompt_file)])
                with open(prompt_file, 'r', encoding='utf-8') as f:
                    prompt = f.read()
            
            else:
                # Start with blank
                prompt = f"Create a {duration}-minute podcast script about '{topic}'."
                save_prompt(prompt, project_name, "blank_prompt.txt")
                prompt_file = Path(f"./projects/{project_name}/prompts/blank_prompt.txt")
                subprocess.run([get_text_editor(), str(prompt_file)])
                with open(prompt_file, 'r', encoding='utf-8') as f:
                    prompt = f.read()
    
            prompt = f"""{prompt}

        === RESEARCH CONTEXT AND INSTRUCTIONS ===

        {research_context}

        IMPORTANT: Follow the research instructions above. Conduct thorough online research using web search. Find and analyze the specified number of sources. Document your sources at the end of the script."""
    
        # 7a. Check and process source documents BEFORE prompt review
        source_documents = process_source_documents(project_name)
        if source_documents:
            prompt = f"""{prompt}

    === USER-PROVIDED SOURCE DOCUMENTS ===

    The following documents were provided by the user. Reference and cite them where relevant alongside your web research:

    {source_documents}

    ===================================
    """
            print(f"\n[INFO] Added {len(source_documents)} characters from source documents to prompt")
    
        # 7b. Count documents for estimation
        doc_count = 0
        if source_documents:
            doc_count = len([s for s in source_documents.split('### SOURCE:') if s.strip()])

        # 7c. Multi-call mode: Show generation plan and get confirmation
        gen_config = config.get('script_generation', {})
        use_multi_call = gen_config.get('enable_multi_call', False) and not is_test_mode

        if use_multi_call:
            if not display_generation_plan(duration, doc_count, web_source_count, config):
                print("Cancelled")
                return

        # 7d. Review final prompt (including source documents if added)
        print("\n" + "="*60)
        print("PROMPT REVIEW")
        print("="*60)
        print(f"Topic: {topic}")
        print(f"Duration: {duration} minutes (~{word_count} words)")
        print(f"Style: {config['styles'][selected_style]['description']}")
        print(f"Language: {config['languages'][selected_language]['name']}")
        if use_multi_call:
            print(f"Web Sources: {web_source_count} (multi-call research)")
        if doc_count > 0:
            print(f"Source Documents: {doc_count} document(s) attached")
        else:
            print(f"Source Documents: None")
        if use_multi_call:
            print(f"Mode: MULTI-CALL (scalable generation)")
        else:
            print(f"Mode: SINGLE-CALL (legacy)")
        print("="*60)

        if not use_multi_call:
            # Legacy mode: show prompt file location
            print("\nFull prompt saved for your review if needed.")
            print(f"Location: {Path(f'./projects/{project_name}/prompts/temp_prompt.txt').absolute()}")
            print("="*60)

            temp_prompt_path = Path(f"./projects/{project_name}/prompts/temp_prompt.txt")
            save_prompt(prompt, project_name, "temp_prompt.txt")

            confirm = get_user_input("\nOptions", [
                "Confirm and send to Claude",
                "Edit prompt in text editor",
                "Cancel"
            ])

            if confirm == 1:
                print(f"\nOpening prompt in your text editor...")
                subprocess.run([get_text_editor(), str(temp_prompt_path)])
                with open(temp_prompt_path, 'r', encoding='utf-8') as f:
                    prompt = f.read()
                print("✓ Prompt updated")
            elif confirm == 2:
                print("Cancelled")
                return
        else:
            # Multi-call mode: already confirmed in generation plan
            pass

        print("\n" + "="*60)
        print("STARTING SCRIPT GENERATION...")
        print("="*60)

        # 8. Generate script
        if use_multi_call:
            # Multi-call generation
            style_description = config['styles'][selected_style]['description']

            # Get style template for reference
            template_file = config['styles'][selected_style]['default_template_file']
            template_file = template_file.replace('{language}', selected_language)
            style_template = ""
            if Path(template_file).exists():
                with open(template_file, 'r', encoding='utf-8') as f:
                    style_template = f.read()
                # Substitute provider-specific placeholders
                style_template = substitute_template_placeholders(
                    style_template, selected_provider, duration
                )

            script = run_multi_call_generation(
                topic=topic,
                duration=duration,
                word_count=word_count,
                research_context=research_context,
                source_documents=source_documents,
                web_source_count=web_source_count,
                style_template=style_template,
                style_description=style_description,
                language=selected_language,
                api_key=anthropic_key,
                config=config,
                project_name=project_name,
                provider=selected_provider
            )

            if not script:
                print("Failed to generate script")
                return

        else:
            # Legacy single-call generation
            script, claude_usage = generate_script(prompt, anthropic_key)
            if not script:
                print("Failed to generate script")
                return

            script = extract_and_save_sources(script, project_name)
    
        draft_num = 1
        # Use provider tag (CRTS/11LB) in script filename
        script_tag = provider_tag
        if is_test_mode:
            script_path = save_script_test(script, project_name, language_code, topic_tag, script_tag, draft_num)
        else:
            script_path = save_script(script, project_name, language_code, script_tag, draft_num)
        print(f"Script generated! ({len(script.split())} words)")
        print(f"Saved to: {script_path}")
    
        # 9. Review and revision loop
        while True:
            print("\n" + "="*60)
            print("SCRIPT REVIEW")
            print("="*60)
            print(f"Script location: {script_path}")
            print("="*60)
        
            action = get_user_input("\nWhat would you like to do?", [
                "Open script in text editor to review",
                "Approve script and proceed to audio",
                "Ask Claude to revise (provide guidance)",
                "Edit script file manually, then regenerate from edits",
                "Save prompt variant to project",
                "Cancel"
            ])
        
            if action == 0:
                print(f"\nOpening {script_path} in your text editor...")
                subprocess.run([get_text_editor(), str(script_path)])
                print("\n✓ Editor closed")
            
            elif action == 1:
                if not validate_template_quality(script):
                    continue
                break
            
            elif action == 2:
                print("\nProvide specific guidance for what to change.")
                print("Examples:")
                print("  - Add more interruptions and overlapping dialogue")
                print("  - Make it more casual - use 'Du' form and colloquialisms")
                print("  - Add more emotional reactions ([excited], [laughs], etc.)")
                guidance = input("\nRevision guidance: ")
            
                if not guidance.strip():
                    print("No guidance provided, skipping revision")
                    continue
                
                revised = revise_script(script, guidance, anthropic_key)
                if revised:
                    script = extract_and_save_sources(revised, project_name)
                    draft_num += 1
                    if is_test_mode:
                        script_path = save_script_test(script, project_name, language_code, topic_tag, script_tag, draft_num)
                    else:
                        script_path = save_script(script, project_name, language_code, script_tag, draft_num)
                    print(f"✓ Revised script saved to: {script_path}")
                else:
                    print("✗ Revision failed")
                
            elif action == 3:
                print(f"\n1. Edit {script_path} in your text editor")
                print("2. Save your changes")
                print("3. Come back here and we'll regenerate with Claude")
                input("\nPress Enter when you're ready to regenerate...")
            
                with open(script_path, 'r', encoding='utf-8') as f:
                    edited_script = f.read()
            
                print("\nWhat changes did you make? (This helps Claude understand context)")
                context = input("Your changes: ")
            
                regenerate_prompt = f"""I have a podcast script that was manually edited. Please review it and generate an improved version that:
    1. Maintains the edits and improvements that were made
    2. Ensures consistent dialogue format with Speaker A: and Speaker B: labels
    3. Improves any rough transitions or formatting issues
    4. Keeps the same overall structure and content

    User's notes on their edits: {context}

    Here is the edited script:

    {edited_script}

    Please provide the improved script maintaining all manual edits and improvements."""
            
                regenerated = generate_script(regenerate_prompt, anthropic_key)
                if regenerated:
                    script = extract_and_save_sources(regenerated, project_name)
                    draft_num += 1
                    if is_test_mode:
                        script_path = save_script_test(script, project_name, language_code, topic_tag, script_tag, draft_num)
                    else:
                        script_path = save_script(script, project_name, language_code, script_tag, draft_num)
                    print(f"✓ Regenerated script saved to: {script_path}")
                else:
                    print("✗ Regeneration failed")
                
            elif action == 4:
                filename = input("Enter filename for prompt variant: ")
                if not filename.endswith('.txt'):
                    filename += '.txt'
                save_prompt(prompt, project_name, filename)
                print(f"✓ Prompt saved to: ./projects/{project_name}/prompts/{filename}")
            
            else:
                print("Cancelled")
                return

    # 10. TTS CONFIGURATION (mode and speed - provider already selected)
    print("\n" + "="*60)
    print("AUDIO CONFIGURATION")
    print("="*60)
    print(f"Provider: {selected_provider.upper()} (selected before script generation)")

    # Mode selection
    if selected_provider == 'elevenlabs':
        print("\n[INFO] ElevenLabs supports quality tiers:")
        print("  - Prototype: 64kbps (lower cost, testing)")
        print("  - Production: 128kbps+ (full quality)")
    elif selected_provider == 'cartesia':
        print("\n[INFO] Cartesia note: Always generates full quality")
        print("  (API does not support quality tiers)")

    mode_idx = get_user_input("\nSelect mode", [
        "Prototype (lower quality, reduced cost for testing)",
        "Production (full quality)"
    ])
    mode = "prototype" if mode_idx == 0 else "production"

    # Get speed setting
    default_speed = config['languages'][selected_language]['speed']
    speed_input = input(f"\nSpeech speed (0.7-1.2, default {default_speed}, Enter to use default): ").strip()
    if speed_input:
        try:
            speed = float(speed_input)
            speed = max(0.7, min(1.2, speed))
            print(f"Using speed: {speed}")
        except ValueError:
            speed = default_speed
            print(f"Invalid, using default: {speed}")
    else:
        speed = default_speed
        print(f"Using default speed: {speed}")

    # 11. Generate audio
    print("\n" + "="*60)
    print(f"AUDIO GENERATION - {mode.upper()} MODE")
    print("="*60)
    settings = config['elevenlabs_settings'][mode]
    print(f"Quality: {settings['quality']}")
    if settings.get('downsample_enabled'):
        print(f"Bitrate: {settings['downsample_bitrate']} (downsampled)")
    print("\n[INFO] Enhanced error logging enabled")
    print("[INFO] Debug chunks will be saved to: projects/{}/debug/".format(project_name))
    print("="*60)
    
    confirm = input("\nProceed with audio generation? (Y/n): ")
    if confirm.lower() == 'n':
        print("Cancelled")
        return
    
    script_for_audio = clean_script_for_audio(script)
    
    # FINAL SAFETY CHECK - Guarantee sources not in audio
    print("\n[FINAL CHECK] Verifying cleaned script...")
    if 'SOURCES FOUND' in script_for_audio.upper():
        print("[ERROR] ❌ SOURCES STILL IN SCRIPT!")
        print("Attempting emergency removal...")
        # Emergency fallback - find any variant and cut
        idx = script_for_audio.upper().find('SOURCES FOUND')
        if idx > 0:
            script_for_audio = script_for_audio[:idx]
            print(f"[INFO] Emergency cut at position {idx}")
    
    if re.search(r'\n\d+\.\s+\*\*', script_for_audio):
        print("[WARNING] ⚠️ Numbered list detected - may be sources!")
        print("First 200 chars of end of script:")
        print(script_for_audio[-200:])
        print("\nLast 50 chars:")
        print(repr(script_for_audio[-50:]))
        confirm = input("\nSources may be in audio. Continue anyway? (y/N): ")
        if confirm.lower() != 'y':
            print("Aborted - fix script manually")
            return
    else:
        print("[INFO] ✓ Verified clean - no sources detected")
    
    audio_data, elevenlabs_chars = generate_audio(script_for_audio, config, language_code, selected_provider, mode, speed, project_name)
    if not audio_data:
        print("\n" + "="*60)
        print("AUDIO GENERATION FAILED")
        print("="*60)
        print("Debug information saved to:")
        print(f"  projects/{project_name}/debug/chunk_*_content.json")
        print("\nCheck the debug output above for:")
        print("  - Error type (Timeout, 500, etc.)")
        print("  - Which chunk failed")
        print("  - Response body details")
        print("\nNext steps:")
        print("1. Review the error message")
        print("2. Check failed chunk JSON")
        print("3. Try running again (may be temporary)")
        print("="*60)
        return
    
    try:
        if is_test_mode:
            audio_path = save_audio(audio_data, project_name, topic, language_code, provider_tag, mode, speed, config, is_test_mode=True, topic_tag=topic_tag)
        else:
            audio_path = save_audio(audio_data, project_name, topic, language_code, provider_tag, mode, speed, config)
        print(f"[DEBUG] Audio saved to: {audio_path}")
    except Exception as e:
        print(f"[ERROR] Failed to save audio: {e}")
        print(f"[DEBUG] Attempted path: projects/{project_name}/audio/")
        return

    # Cleanup debug files after successful generation
    debug_dir = Path(f"./projects/{project_name}/debug")
    if debug_dir.exists():
        for debug_file in debug_dir.glob("chunk_*_content.json"):
            debug_file.unlink()
        print("[INFO] ✓ Cleaned up debug files")
    
    # 11. Display results
    print("\n" + "="*60)
    print("✓ PODCAST GENERATED SUCCESSFULLY!")
    print("="*60)
    print(f"File: {audio_path}")
    print(f"Size: {len(audio_data) / 1024 / 1024:.1f} MB")
    print(f"Mode: {mode.upper()}")
    print("="*60)
    
    # 12. Save prompt (only if new script was generated)
    if not script_ready:  # Script was generated, not loaded
        save = input("\nSave final prompt to project? (Y/n): ")
        if save.lower() != 'n':
            filename = input("Enter filename: ")
            if not filename.endswith('.txt'):
                filename += '.txt'
            save_prompt(prompt, project_name, filename)
            print(f"Prompt saved to: ./projects/{project_name}/prompts/{filename}")
    
    # 13. Project summary
    print("\n" + "="*60)
    print("PROJECT SUMMARY")
    print("="*60)
    print(f"Project: {project_name}")
    print(f"Location: ./projects/{project_name}/")
    
    script_count = len(list(Path(f"./projects/{project_name}/scripts/").glob("*draft*.txt")))
    audio_count = len(list(Path(f"./projects/{project_name}/audio/").glob("*.mp3")))
    prompt_count = len(list(Path(f"./projects/{project_name}/prompts/").glob("*.txt")))
    
    print(f"- Scripts: {script_count} drafts")
    print(f"- Audio: {audio_count} files")
    print(f"- Prompts: {prompt_count} files")
    print("="*60)
    
    # 14. Generate another?
    another = input("\nGenerate another podcast? (Y/n): ")
    if another.lower() != 'n':
        print("\n" + "="*60 + "\n")
        main()
    else:
        print("\nPipeline complete!")


if __name__ == "__main__":
    main()
